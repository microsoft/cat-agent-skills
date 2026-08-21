#!/usr/bin/env python3
"""
verify_pptx.py — Hard gate: a PPTX is "good" only if it actually opens.

Verification levels
-------------------
Two levels exist, and the tool always reports which one it achieved:

  render      structural + package load + a real LibreOffice conversion.
              This is the only level that proves the file opens.

  structural  structural + package load, no conversion. A strong signal, but
              NOT proof that PowerPoint will open the file.

Render is REQUIRED by default. If LibreOffice is not present, verification
FAILS with a non-zero exit — it does not quietly downgrade and report success,
because a pass that cannot distinguish "opens" from "was never opened" is worse
than no check at all.

A caller who knowingly accepts the weaker guarantee must ask for it explicitly
with --allow-no-render. That run exits 0 on success but reports
`verified_by: structural`, and the human-readable output says in plain words
that the file is not render-verified. Callers gating a delivery should require
`verified_by == "render"` rather than just `ok == true`.

Note that `soffice` is not guaranteed to exist in the Copilot Studio container.
Install it in the image if you need the render guarantee there; otherwise run
with --allow-no-render and accept that the file has not been proven to open.

Usage:
    python verify_pptx.py <file.pptx> [--allow-no-render] [--json]

    --allow-no-render   Downgrade to structural verification when LibreOffice
                        is unavailable, instead of failing.
    --require-render    Explicit form of the default; fails if soffice is
                        missing. Kept for callers that want it spelled out.
    --json              Machine-readable result on stdout.

Exit codes:
    0  verification passed at the level reported in `verified_by`
    1  verification failed
    2  bad invocation (file missing, unreadable, or not a ZIP)
"""

import argparse
import json
import posixpath
import shutil
import subprocess
import sys
import tempfile
import zipfile
from pathlib import Path

try:
    from lxml import etree
except ImportError:
    sys.exit("ERROR: lxml is required. Install it with: pip install lxml")

_SAFE_PARSER = etree.XMLParser(resolve_entities=False, no_network=True, load_dtd=False)

NS_CT  = "http://schemas.openxmlformats.org/package/2006/content-types"
NS_P   = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS_R   = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
NS_REL = "http://schemas.openxmlformats.org/package/2006/relationships"


def rels_path_for(part: str) -> str:
    d, b = posixpath.split(part)
    return posixpath.join(d, "_rels", b + ".rels")


def part_for_rels(rels_name: str) -> str:
    """'ppt/slides/_rels/slide1.xml.rels' -> 'ppt/slides/slide1.xml'"""
    d, b = posixpath.split(rels_name)
    return posixpath.join(posixpath.dirname(d), b[:-len(".rels")])


# ── structural ──────────────────────────────────────────────────────────────
def structural_checks(path: Path):
    errors, warnings = [], []

    try:
        z = zipfile.ZipFile(path)
    except zipfile.BadZipFile as exc:
        return [f"not a valid ZIP/OPC package — {exc}"], warnings

    with z:
        bad = z.testzip()
        if bad is not None:
            errors.append(f"ZIP CRC error in {bad}")
        names = set(z.namelist())

        defaults, overrides = {}, {}
        if "[Content_Types].xml" not in names:
            errors.append("[Content_Types].xml missing")
        else:
            try:
                root = etree.fromstring(z.read("[Content_Types].xml"), _SAFE_PARSER)
            except etree.XMLSyntaxError as exc:
                errors.append(f"[Content_Types].xml: malformed XML — {exc}")
                root = None
            if root is not None:
                # The Types element MUST be in the default (unprefixed) namespace.
                if root.prefix is not None:
                    errors.append(
                        f"[Content_Types].xml uses namespace prefix '{root.prefix}:' — "
                        "must be the default namespace or the package will not open.")
                for d in root.findall(f"{{{NS_CT}}}Default"):
                    ext = (d.get("Extension") or "").lower()
                    if ext:
                        defaults[ext] = d.get("ContentType")
                for ov in root.findall(f"{{{NS_CT}}}Override"):
                    pn = (ov.get("PartName") or "").lstrip("/")
                    if pn:
                        overrides[pn] = ov.get("ContentType")
                    if pn and pn not in names:
                        errors.append(f"[Content_Types].xml: Override for missing part '/{pn}'")

        # Every part needs a declared content type.
        for n in sorted(names):
            if n == "[Content_Types].xml" or n.endswith("/"):
                continue
            base = posixpath.basename(n)
            ext = base.rsplit(".", 1)[-1].lower() if "." in base else ""
            if n not in overrides and ext not in defaults:
                errors.append(f"{n}: no content type declared (no Override, no Default '{ext}')")

        # Every internal relationship Target must resolve to a part that exists,
        # and every r:* reference in a part must resolve in that part's .rels.
        for n in sorted(names):
            if not n.endswith(".rels"):
                continue
            try:
                root = etree.fromstring(z.read(n), _SAFE_PARSER)
            except etree.XMLSyntaxError as exc:
                errors.append(f"{n}: malformed XML — {exc}")
                continue
            if root.prefix is not None:
                errors.append(
                    f"{n} uses namespace prefix '{root.prefix}:' — "
                    "must be the default namespace or the package will not open.")
            source_part = part_for_rels(n)
            part_folder = posixpath.dirname(source_part)
            seen_ids = set()
            for rel in root:
                rid = rel.get("Id", "")
                if rid in seen_ids:
                    errors.append(f"{n}: duplicate relationship Id '{rid}'")
                seen_ids.add(rid)
                if rel.get("TargetMode", "") == "External":
                    continue
                tgt = rel.get("Target", "")
                if tgt.startswith("/"):
                    errors.append(f"{n}: absolute internal Target '{tgt}' "
                                  "(needs relative path or TargetMode=External)")
                    continue
                resolved = posixpath.normpath(posixpath.join(part_folder, tgt))
                if resolved not in names:
                    errors.append(f"{n}: Target '{tgt}' resolves to "
                                  f"'{resolved}' which is missing from the package")

            if source_part in names and source_part.endswith(".xml"):
                try:
                    part_root = etree.fromstring(z.read(source_part), _SAFE_PARSER)
                except etree.XMLSyntaxError as exc:
                    errors.append(f"{source_part}: malformed XML — {exc}")
                    continue
                for el in part_root.iter():
                    for k, v in el.attrib.items():
                        if not k.startswith("{" + NS_R + "}"):
                            continue
                        if isinstance(v, str) and v.startswith("rId") and v not in seen_ids:
                            errors.append(
                                f"{source_part}: @{etree.QName(k).localname}='{v}' "
                                f"has no matching relationship in {n}")

        # presentation.xml sanity: masters and slides have required ids.
        if "ppt/presentation.xml" not in names:
            errors.append("ppt/presentation.xml missing")
        else:
            try:
                pr = etree.fromstring(z.read("ppt/presentation.xml"), _SAFE_PARSER)
            except etree.XMLSyntaxError as exc:
                errors.append(f"ppt/presentation.xml: malformed XML — {exc}")
                pr = None
            if pr is not None:
                gids = []
                ml = pr.find(f".//{{{NS_P}}}sldMasterIdLst")
                if ml is None or len(ml) == 0:
                    errors.append("presentation.xml has no slide masters")
                else:
                    for e in ml:
                        if not e.get("id"):
                            errors.append("sldMasterId missing required 'id' attribute")
                        else:
                            gids.append(e.get("id"))
                sl = pr.find(f".//{{{NS_P}}}sldIdLst")
                if sl is None or len(sl) == 0:
                    warnings.append("presentation.xml lists no slides")
                for n in sorted(names):
                    if n.startswith("ppt/slideMasters/slideMaster") and n.endswith(".xml"):
                        m = etree.fromstring(z.read(n), _SAFE_PARSER)
                        ll = m.find(f".//{{{NS_P}}}sldLayoutIdLst")
                        if ll is not None:
                            for e in ll:
                                if e.get("id"):
                                    gids.append(e.get("id"))
                dupes = {x for x in gids if gids.count(x) > 1}
                if dupes:
                    errors.append(f"Duplicate global master/layout IDs: {sorted(dupes)}")

    return errors, warnings


# ── package load ────────────────────────────────────────────────────────────
def load_check(path: Path):
    """Open the package with a real OPC consumer.

    This is not a render, but it is a genuine load: python-pptx resolves the
    content types, walks the relationship graph and instantiates every slide,
    layout and master, so most packaging faults surface here. Pure Python, so
    it is available wherever the merge scripts themselves run.
    """
    try:
        from pptx import Presentation
    except ImportError:
        return None, "python-pptx not installed; load check skipped"
    try:
        prs = Presentation(str(path))
        n_slides = len(prs.slides._sldIdLst)
        n_masters = len(prs.slide_masters)
        n_layouts = sum(len(m.slide_layouts) for m in prs.slide_masters)
        for slide in prs.slides:
            _ = [sh.shape_type for sh in slide.shapes]
        return True, (f"loaded {n_slides} slides, {n_masters} masters, "
                      f"{n_layouts} layouts")
    except Exception as exc:                     # noqa: BLE001 — report anything
        return False, f"{type(exc).__name__}: {exc}"


# ── render ──────────────────────────────────────────────────────────────────
def find_soffice():
    return shutil.which("soffice") or shutil.which("libreoffice")


def render_check(path: Path, soffice: str):
    with tempfile.TemporaryDirectory() as td:
        prof = Path(td) / "profile"
        outdir = Path(td) / "out"
        outdir.mkdir()
        try:
            proc = subprocess.run(
                [soffice, "--headless", f"-env:UserInstallation=file://{prof}",
                 "--convert-to", "pdf", "--outdir", str(outdir), str(path)],
                capture_output=True, timeout=180,
            )
        except subprocess.TimeoutExpired:
            return False, "render timed out after 180s"
        pdfs = list(outdir.glob("*.pdf"))
        if pdfs and pdfs[0].stat().st_size > 0:
            return True, f"rendered to {pdfs[0].stat().st_size} byte PDF"
        detail = (proc.stderr or proc.stdout or b"").decode("utf-8", "replace").strip()
        return False, f"LibreOffice could not load the file{': ' + detail if detail else ''}"


# ── main ────────────────────────────────────────────────────────────────────
def main() -> int:
    ap = argparse.ArgumentParser(
        description="Verify a PPTX. Render is required unless --allow-no-render.")
    ap.add_argument("pptx")
    ap.add_argument("--allow-no-render", action="store_true",
                    help="Downgrade to structural verification when LibreOffice "
                         "is unavailable, instead of failing.")
    ap.add_argument("--require-render", action="store_true",
                    help="Explicit form of the default behaviour.")
    ap.add_argument("--json", action="store_true")
    args = ap.parse_args()

    path = Path(args.pptx)
    result = {
        "ok": False,
        "stage": "validate",
        "file": str(path),
        "verified_by": None,          # "render" | "structural" | None
        "render_required": not args.allow_no_render,
        "structural_errors": [],
        "warnings": [],
        "load": None,
        "render": None,
    }

    def emit(code):
        if args.json:
            print(json.dumps(result, indent=2))
        else:
            if result["ok"]:
                if result["verified_by"] == "render":
                    print(f"PASS (render-verified) — {result['render']['detail']}")
                else:
                    print("PASS (structural only) — NOT render-verified; "
                          "this file has not been proven to open.")
                    print(f"  load: {result['load']['detail']}")
            else:
                print("FAIL")
                # load and render failures are already folded into
                # structural_errors, so this list is the complete picture.
                for e in result["structural_errors"]:
                    print(f"  - {e}")
            for w in result["warnings"]:
                print(f"  ! {w}")
        return code

    if not path.exists():
        result["structural_errors"] = ["file does not exist"]
        return emit(2)

    errors, warnings = structural_checks(path)
    result["structural_errors"] = list(errors)
    result["warnings"] = warnings
    if errors and errors[0].startswith("not a valid ZIP"):
        return emit(2)

    load_ok, load_msg = load_check(path)
    result["load"] = {"passed": load_ok, "detail": load_msg}
    if load_ok is False:
        errors.append(f"load: {load_msg}")
    elif load_ok is None:
        warnings.append(load_msg)

    soffice = find_soffice()
    if soffice:
        render_ok, render_msg = render_check(path, soffice)
        result["render"] = {"passed": render_ok, "detail": render_msg}
        if render_ok is False:
            errors.append(f"render: {render_msg}")
    else:
        msg = ("LibreOffice (soffice/libreoffice) not found — the file cannot be "
               "proven to open in this environment")
        if args.allow_no_render:
            result["render"] = {"passed": None, "detail": msg + " (accepted via --allow-no-render)"}
            warnings.append("render check skipped: " + msg)
        else:
            result["render"] = {"passed": False, "detail": msg}
            errors.append(
                "render: " + msg + ". Install LibreOffice in the container, or "
                "re-run with --allow-no-render to accept structural verification only.")

    result["structural_errors"] = errors
    result["ok"] = not errors
    if result["ok"]:
        result["verified_by"] = "render" if (result["render"]
                                             and result["render"]["passed"]) else "structural"
    return emit(0 if result["ok"] else 1)


if __name__ == "__main__":
    sys.exit(main())
