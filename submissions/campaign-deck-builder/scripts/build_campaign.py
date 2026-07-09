#!/usr/bin/env python3
"""Build a 3-slide marketing-campaign deck from a JSON answers file.

Fills the bundled ``assets/campaign-template.pptx`` — replacing every
``{{TOKEN}}`` with your answers (preserving the template's formatting) and
optionally recoloring the deck to your brand accent color.

Usage
-----
    python scripts/build_campaign.py answers.json
    python scripts/build_campaign.py answers.json campaign-deck.pptx
    python scripts/build_campaign.py answers.json out.pptx --template other.pptx

The answers file is JSON. All keys are optional — anything you omit falls back
to a tasteful default, so a partial file still produces a complete deck::

    {
      "brand_name": "Vela",
      "campaign_name": "Taste the Quiet",
      "tagline": "Botanical sparkling water for your loudest days.",
      "audience": "calm-seeking millennials",
      "big_idea": "A moment of quiet in every can.",
      "messages": [
        "Zero sugar, all serenity",
        "Adaptogens that actually work",
        "100% recyclable, guilt-free fizz"
      ],
      "channels": "TikTok  ·  Out-of-home  ·  Sampling  ·  Creator gifting",
      "cta": "Pop the quiet.",
      "closing": "Vela — find your fizz.",
      "accent_hex": "#00A6A6"
    }
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
except ImportError:  # pragma: no cover
    sys.exit(
        "python-pptx is required. Install it with:  pip install python-pptx"
    )

TOKEN_RE = re.compile(r"\{\{([A-Z0-9_]+)\}\}")

# Sensible, on-theme defaults so a partial (or empty) answers file still builds.
DEFAULTS = {
    "brand_name": "Your Brand",
    "campaign_name": "Campaign Name",
    "tagline": "A tagline that makes people lean in.",
    "audience": "your target audience",
    "big_idea": "The one bold idea this whole campaign turns on.",
    "messages": [
        "Key message one",
        "Key message two",
        "Key message three",
    ],
    "channels": "Social  ·  Out-of-home  ·  Influencer  ·  Retail",
    "cta": "Here's the ask.",
    "closing": "One line to land it.",
    "accent_hex": "",  # empty => keep the template's vibrant default accent
}


def load_answers(path: Path) -> dict:
    try:
        raw = json.loads(path.read_text(encoding="utf-8"))
    except FileNotFoundError:
        sys.exit(f"answers file not found: {path}")
    except json.JSONDecodeError as exc:
        sys.exit(f"answers file is not valid JSON: {exc}")
    if not isinstance(raw, dict):
        sys.exit("answers file must be a JSON object of field -> value")
    return raw


def normalize_hex(value: str) -> str | None:
    """Return an uppercase 6-digit hex (no '#') or None if invalid/empty."""
    if not value:
        return None
    v = value.strip().lstrip("#")
    if re.fullmatch(r"[0-9a-fA-F]{6}", v):
        return v.upper()
    if re.fullmatch(r"[0-9a-fA-F]{3}", v):  # allow shorthand #abc
        return "".join(c * 2 for c in v).upper()
    return None


def build_token_map(answers: dict) -> dict[str, str]:
    """Merge answers over defaults and flatten into a {{TOKEN}} -> text map."""

    def pick(key):
        val = answers.get(key)
        if val is None or (isinstance(val, str) and not val.strip()):
            return DEFAULTS[key]
        return val

    # messages can be a list, or discrete message_1/2/3 keys.
    messages = answers.get("messages")
    if not isinstance(messages, list) or not messages:
        messages = [
            answers.get(f"message_{i}") or DEFAULTS["messages"][i - 1]
            for i in range(1, 4)
        ]
    messages = (list(messages) + DEFAULTS["messages"])[:3]

    return {
        "BRAND_NAME": str(pick("brand_name")),
        "CAMPAIGN_NAME": str(pick("campaign_name")),
        "TAGLINE": str(pick("tagline")),
        "AUDIENCE": str(pick("audience")),
        "BIG_IDEA": str(pick("big_idea")),
        "MESSAGE_1": str(messages[0]),
        "MESSAGE_2": str(messages[1]),
        "MESSAGE_3": str(messages[2]),
        "CHANNELS": str(pick("channels")),
        "CTA": str(pick("cta")),
        "CLOSING": str(pick("closing")),
    }


def _sub(text: str, token_map: dict[str, str]) -> str:
    return TOKEN_RE.sub(lambda m: token_map.get(m.group(1), m.group(0)), text)


def replace_tokens(slide, token_map: dict[str, str]) -> None:
    """Replace {{TOKEN}}s while preserving each run's formatting.

    The template authors every token in its own run, so a per-run pass keeps
    all formatting intact. A paragraph-level fallback handles the rare case of a
    token split across runs (e.g. after manual template edits)."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                if "{{" in run.text:
                    run.text = _sub(run.text, token_map)
            # Fallback: token split across runs -> collapse into the first run.
            joined = "".join(r.text for r in para.runs)
            if "{{" in joined and TOKEN_RE.search(joined):
                new_text = _sub(joined, token_map)
                if para.runs:
                    para.runs[0].text = new_text
                    for extra in para.runs[1:]:
                        extra.text = ""


def apply_accent(prs, hex6: str | None, warnings: list[str]) -> int:
    """Recolor every shape named 'accent*' to the brand hex. Returns count."""
    if not hex6:
        return 0
    color = RGBColor.from_string(hex6)
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.name.lower().startswith("accent"):
                continue
            try:
                shape.fill.solid()
                shape.fill.fore_color.rgb = color
                count += 1
            except (AttributeError, TypeError):  # shape can't be solid-filled
                warnings.append(f"could not recolor shape '{shape.name}'")
    return count


def main(argv=None) -> int:
    parser = argparse.ArgumentParser(
        description="Fill the campaign .pptx template from a JSON answers file."
    )
    parser.add_argument("answers", type=Path, help="path to the JSON answers file")
    parser.add_argument(
        "output",
        type=Path,
        nargs="?",
        default=Path("campaign-deck.pptx"),
        help="output .pptx path (default: campaign-deck.pptx)",
    )
    parser.add_argument(
        "--template",
        type=Path,
        default=None,
        help="template .pptx (default: assets/campaign-template.pptx next to this script)",
    )
    args = parser.parse_args(argv)

    template = args.template or (
        Path(__file__).resolve().parent.parent / "assets" / "campaign-template.pptx"
    )
    if not template.exists():
        sys.exit(f"template not found: {template}")

    answers = load_answers(args.answers)
    token_map = build_token_map(answers)

    warnings: list[str] = []
    raw_hex = str(answers.get("accent_hex", "") or "")
    hex6 = normalize_hex(raw_hex)
    if raw_hex and hex6 is None:
        warnings.append(
            f"ignoring invalid accent_hex '{raw_hex}' (want #RRGGBB); "
            "keeping template default"
        )

    prs = Presentation(str(template))
    for slide in prs.slides:
        replace_tokens(slide, token_map)
    recolored = apply_accent(prs, hex6, warnings)

    # Report any tokens that survived (typo in a template token name, etc.).
    leftover = set()
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                leftover.update(TOKEN_RE.findall(shape.text_frame.text))
    for tok in sorted(leftover):
        warnings.append(f"unfilled placeholder left in deck: {{{{{tok}}}}}")

    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.output))

    print(f"✓ Built {args.output}")
    print(f"  brand   : {token_map['BRAND_NAME']}")
    print(f"  campaign: {token_map['CAMPAIGN_NAME']}")
    print(f"  slides  : {len(prs.slides._sldIdLst)}")
    if hex6:
        print(f"  accent  : #{hex6} ({recolored} shapes recolored)")
    for w in warnings:
        print(f"  ! {w}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
