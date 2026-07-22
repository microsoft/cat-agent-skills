#!/usr/bin/env python3
"""Inventory a document library and generate a human-review curation backlog."""

from __future__ import annotations

import argparse
import csv
import hashlib
import html
import json
import math
import re
import sys
from collections import Counter, defaultdict
from datetime import datetime, timezone
from pathlib import Path
from typing import Any, Iterable

SUPPORTED_EXTENSIONS = {
    ".csv",
    ".docx",
    ".htm",
    ".html",
    ".jpeg",
    ".jpg",
    ".json",
    ".md",
    ".pdf",
    ".png",
    ".pptx",
    ".txt",
    ".xlsx",
    ".xml",
    ".yaml",
    ".yml",
}
DETECTED_TYPE_TO_EXTENSION = {
    "csv": ".csv",
    "docx": ".docx",
    "html": ".html",
    "jpeg": ".jpg",
    "json": ".json",
    "markdown": ".md",
    "pdf": ".pdf",
    "png": ".png",
    "pptx": ".pptx",
    "text": ".txt",
    "txt": ".txt",
    "xlsx": ".xlsx",
    "xml": ".xml",
    "yaml": ".yaml",
}
WORD_RE = re.compile(r"[A-Za-z0-9][A-Za-z0-9'_-]*")
SENTENCE_RE = re.compile(r"(?<=[.!?])\s+|\n+")
DATE_RE = re.compile(
    r"\b(?:20\d{2}[-/]\d{1,2}[-/]\d{1,2}|"
    r"\d{1,2}[-/]\d{1,2}[-/]20\d{2})\b"
)
NUMBER_RE = re.compile(r"\b\d+(?:\.\d+)?%?\b")
NEGATIVE_DIRECTIVES = ("must not", "should not", "prohibited", "never")
POSITIVE_DIRECTIVES = ("must", "should", "required", "always")
PRIORITY_ORDER = {"Critical": 0, "High": 1, "Medium": 2, "Low": 3}
CATEGORY_ORDER = {
    "Potential conflict": 0,
    "Extraction gap": 1,
    "Duplicate": 2,
    "Near duplicate": 3,
    "Stale candidate": 4,
    "Related content": 5,
}
BACKLOG_COLUMNS = (
    "priority",
    "category",
    "recommendedAction",
    "primaryDocument",
    "primaryPath",
    "primaryPage",
    "primaryExcerpt",
    "relatedDocument",
    "relatedPath",
    "relatedPage",
    "relatedExcerpt",
    "confidence",
    "reason",
    "sourceUrl",
    "owner",
)
DOCUMENT_COLUMNS = (
    "path",
    "name",
    "extension",
    "detectedType",
    "detectedMimeType",
    "effectiveExtension",
    "fileTypeMismatch",
    "bytes",
    "modifiedAt",
    "sourceUrl",
    "owner",
    "title",
    "contentType",
    "sha256",
    "normalizedSha256",
    "characters",
    "words",
    "pageCount",
    "extractionStatus",
    "extractionMethod",
    "error",
)
REPORT_COLUMNS = (
    ("priority", "Priority"),
    ("category", "Category"),
    ("primaryDocument", "Primary document"),
    ("primaryPage", "Primary page"),
    ("primaryExcerpt", "Primary excerpt"),
    ("relatedDocument", "Related document"),
    ("relatedPage", "Related page"),
    ("relatedExcerpt", "Related excerpt"),
    ("confidence", "Confidence"),
    ("recommendedAction", "Recommended action"),
    ("reason", "Reason"),
)
TOPIC_STOPWORDS = {
    "a",
    "an",
    "and",
    "are",
    "be",
    "for",
    "from",
    "in",
    "is",
    "of",
    "on",
    "or",
    "the",
    "to",
    "with",
    "must",
    "not",
    "should",
    "required",
    "prohibited",
    "never",
    "always",
    "optional",
}


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input", required=True, type=Path, help="Library directory")
    parser.add_argument(
        "--output",
        type=Path,
        default=Path("/app/created/knowledge-corpus-curation"),
        help="Output directory (default: /app/created/knowledge-corpus-curation)",
    )
    parser.add_argument("--config", type=Path, help="Configuration JSON")
    parser.add_argument("--metadata", type=Path, help="Optional source metadata JSON")
    parser.add_argument(
        "--batch-manifest",
        type=Path,
        help="Manifest produced by prepare_batches.py for source ZIP traceability",
    )
    parser.add_argument("--stale-after-days", type=int, help="Override stale threshold")
    parser.add_argument(
        "--corpus-scope",
        choices=("whole-library", "subset"),
        help="User-confirmed coverage of the uploaded corpus",
    )
    parser.add_argument(
        "--content-scope",
        choices=("current-only", "include-drafts-and-history"),
        help="User-confirmed treatment of drafts, archives, and historical versions",
    )
    parser.add_argument("--ocr", action="store_true", help="OCR images and scanned PDFs")
    parser.add_argument(
        "--no-embeddings", action="store_true", help="Disable local embedding analysis"
    )
    return parser.parse_args()


def load_json(path: Path | None, default: Any) -> Any:
    if path is None:
        return default
    try:
        return json.loads(path.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError) as exc:
        raise SystemExit(f"Cannot read JSON file {path}: {exc}") from exc


def load_metadata(path: Path | None) -> dict[str, dict[str, Any]]:
    raw = load_json(path, {})
    if isinstance(raw, list):
        return {
            normalize_relative_path(str(item.get("path", ""))): item
            for item in raw
            if isinstance(item, dict) and item.get("path")
        }
    if isinstance(raw, dict):
        return {
            normalize_relative_path(str(key)): value
            for key, value in raw.items()
            if isinstance(value, dict)
        }
    raise SystemExit("Metadata JSON must be an array or object.")


def load_batch_sources(path: Path | None) -> list[str]:
    if path is None:
        return []
    raw = load_json(path, {})
    if not isinstance(raw, dict) or not isinstance(raw.get("batches"), list):
        raise SystemExit("Batch manifest must contain a batches array.")
    sources: list[str] = []
    for batch in raw["batches"]:
        if not isinstance(batch, dict) or not batch.get("source"):
            raise SystemExit("Every batch manifest entry must include a source.")
        source = str(batch["source"])
        if source not in sources:
            sources.append(source)
    return sources


def normalize_relative_path(value: str) -> str:
    return value.replace("\\", "/").lstrip("./")


def sha256_bytes(data: bytes) -> str:
    return hashlib.sha256(data).hexdigest()


def normalize_text(text: str) -> str:
    return " ".join(WORD_RE.findall(text.lower()))


def tokenize(text: str) -> list[str]:
    return WORD_RE.findall(text.lower())


def extract_docx(path: Path) -> str:
    from docx import Document

    document = Document(path)
    parts = [paragraph.text for paragraph in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    for section in document.sections:
        parts.extend(p.text for p in section.header.paragraphs)
        parts.extend(p.text for p in section.footer.paragraphs)
    return "\n".join(parts)


def extract_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(path)
    parts: list[str] = []
    for index, slide in enumerate(presentation.slides, start=1):
        parts.append(f"[Slide {index}]")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(cell.text for cell in row.cells))
        try:
            notes = slide.notes_slide.notes_text_frame.text
            if notes:
                parts.append(f"[Notes] {notes}")
        except (AttributeError, KeyError):
            pass
    return "\n".join(parts)


def extract_xlsx(path: Path) -> str:
    from openpyxl import load_workbook

    workbook = load_workbook(path, read_only=True, data_only=True)
    parts: list[str] = []
    try:
        for sheet in workbook.worksheets:
            parts.append(f"[Sheet: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value) for value in row]
                if any(values):
                    parts.append(" | ".join(values))
    finally:
        workbook.close()
    return "\n".join(parts)


def extract_pdf(path: Path, use_ocr: bool) -> tuple[str, str, list[str]]:
    import pdfplumber

    pages: list[str] = []
    with pdfplumber.open(path) as pdf:
        pages = [(page.extract_text() or "") for page in pdf.pages]
    text = "\n".join(pages)
    if text.strip() or not use_ocr:
        return text, "text", pages
    ocr_pages = ocr_pdf(path)
    return "\n".join(ocr_pages), "ocr", ocr_pages


def get_ocr_engine():
    import numpy as np
    from rapidocr_onnxruntime import RapidOCR

    return RapidOCR(), np


def ocr_image(path: Path) -> str:
    from PIL import Image

    engine, np = get_ocr_engine()
    result, _ = engine(np.asarray(Image.open(path).convert("RGB")))
    return "\n".join(item[1] for item in (result or []))


def ocr_pdf(path: Path) -> list[str]:
    import pypdfium2 as pdfium

    engine, np = get_ocr_engine()
    document = pdfium.PdfDocument(path)
    try:
        parts: list[str] = []
        for index in range(len(document)):
            image = document[index].render(scale=2).to_pil().convert("RGB")
            result, _ = engine(np.asarray(image))
            parts.append("\n".join(item[1] for item in (result or [])))
        return parts
    finally:
        document.close()


def extract_html(path: Path) -> str:
    from bs4 import BeautifulSoup

    soup = BeautifulSoup(path.read_text(encoding="utf-8", errors="replace"), "html.parser")
    for node in soup(["script", "style", "noscript"]):
        node.decompose()
    return soup.get_text("\n", strip=True)


def extract_csv(path: Path) -> str:
    parts: list[str] = []
    with path.open("r", encoding="utf-8-sig", errors="replace", newline="") as handle:
        for row in csv.reader(handle):
            parts.append(" | ".join(row))
    return "\n".join(parts)


def extract_text(
    path: Path, use_ocr: bool, effective_extension: str
) -> tuple[str, str, list[str]]:
    extension = effective_extension
    if extension == ".docx":
        return extract_docx(path), "document", []
    if extension == ".pptx":
        return extract_pptx(path), "presentation", []
    if extension == ".xlsx":
        return extract_xlsx(path), "workbook", []
    if extension == ".pdf":
        return extract_pdf(path, use_ocr)
    if extension in {".html", ".htm"}:
        return extract_html(path), "html", []
    if extension == ".csv":
        return extract_csv(path), "table", []
    if extension in {".png", ".jpg", ".jpeg"}:
        if not use_ocr:
            return "", "needs_ocr", []
        return ocr_image(path), "ocr", []
    return path.read_text(encoding="utf-8", errors="replace"), "text", []


def create_file_type_detector(warnings: list[str]):
    try:
        from magika import Magika

        return Magika()
    except Exception as exc:
        warnings.append(
            "Magika unavailable; filename extensions were used as a fallback. "
            f"{type(exc).__name__}: {exc}"
        )
        return None


def detect_file_type(
    path: Path, detector: Any, warnings: list[str]
) -> tuple[str, str]:
    if detector is None:
        return "", ""
    try:
        result = detector.identify_path(path)
        output = result.output
        return str(output.label or "").lower(), str(output.mime_type or "")
    except Exception as exc:
        warnings.append(
            f"{path.name}: Magika detection failed; extension fallback used. "
            f"{type(exc).__name__}: {exc}"
        )
        return "", ""


def parse_datetime(value: Any, fallback: datetime) -> datetime:
    if isinstance(value, str) and value.strip():
        candidate = value.strip().replace("Z", "+00:00")
        try:
            parsed = datetime.fromisoformat(candidate)
            if parsed.tzinfo is None:
                parsed = parsed.replace(tzinfo=timezone.utc)
            return parsed.astimezone(timezone.utc)
        except ValueError:
            pass
    return fallback


def inventory(
    input_root: Path,
    metadata: dict[str, dict[str, Any]],
    minimum_chars: int,
    use_ocr: bool,
) -> tuple[
    list[dict[str, Any]], dict[str, str], dict[str, list[str]], list[str]
]:
    documents: list[dict[str, Any]] = []
    extracted: dict[str, str] = {}
    page_texts: dict[str, list[str]] = {}
    warnings: list[str] = []
    detector = create_file_type_detector(warnings)
    for path in sorted(input_root.rglob("*")):
        if not path.is_file():
            continue
        relative = normalize_relative_path(str(path.relative_to(input_root)))
        source = metadata.get(relative, {})
        raw = path.read_bytes()
        detected_type, detected_mime_type = detect_file_type(path, detector, warnings)
        filename_extension = path.suffix.lower()
        detected_extension = DETECTED_TYPE_TO_EXTENSION.get(detected_type, "")
        effective_extension = detected_extension or filename_extension
        mismatch = bool(
            detected_extension
            and filename_extension
            and detected_extension != filename_extension
            and not (
                detected_extension == ".jpg" and filename_extension == ".jpeg"
            )
        )
        filesystem_modified = datetime.fromtimestamp(path.stat().st_mtime, timezone.utc)
        modified = parse_datetime(source.get("modifiedAt"), filesystem_modified)
        record: dict[str, Any] = {
            "id": f"doc-{len(documents) + 1:04d}",
            "path": relative,
            "name": path.name,
            "extension": path.suffix.lower(),
            "detectedType": detected_type,
            "detectedMimeType": detected_mime_type,
            "effectiveExtension": effective_extension,
            "fileTypeMismatch": mismatch,
            "bytes": len(raw),
            "modifiedAt": modified.isoformat(),
            "sourceUrl": source.get("sourceUrl", ""),
            "owner": source.get("owner", ""),
            "title": source.get("title", path.stem),
            "contentType": source.get("contentType", ""),
            "sha256": sha256_bytes(raw),
            "normalizedSha256": "",
            "characters": 0,
            "words": 0,
            "pageCount": "",
            "extractionStatus": "pending",
            "extractionMethod": "",
            "error": "",
        }
        if effective_extension not in SUPPORTED_EXTENSIONS:
            record["extractionStatus"] = "unsupported"
            record["error"] = (
                f"Unsupported detected file type: {detected_type or filename_extension or 'unknown'}"
            )
            warnings.append(f"{relative}: {record['error']}")
            documents.append(record)
            continue
        if mismatch:
            warnings.append(
                f"{relative}: filename extension {filename_extension} does not "
                f"match detected content type {detected_type}."
            )
        try:
            text, method, pages = extract_text(path, use_ocr, effective_extension)
            normalized = normalize_text(text)
            record["characters"] = len(text)
            record["words"] = len(tokenize(text))
            record["pageCount"] = len(pages) if pages else ""
            record["normalizedSha256"] = (
                hashlib.sha256(normalized.encode("utf-8")).hexdigest()
                if normalized
                else ""
            )
            record["extractionMethod"] = method
            if method == "needs_ocr":
                record["extractionStatus"] = "needs_ocr"
                warnings.append(f"{relative}: image requires --ocr.")
            elif len(normalized) < minimum_chars:
                record["extractionStatus"] = "insufficient_text"
                warnings.append(
                    f"{relative}: extracted only {len(normalized)} normalized characters."
                )
            else:
                record["extractionStatus"] = "ok"
                extracted[record["id"]] = text
                if pages:
                    page_texts[record["id"]] = pages
        except Exception as exc:
            record["extractionStatus"] = "failed"
            record["error"] = f"{type(exc).__name__}: {exc}"
            warnings.append(f"{relative}: extraction failed: {record['error']}")
        documents.append(record)
    return documents, extracted, page_texts, warnings


def exact_duplicate_groups(
    documents: list[dict[str, Any]], field: str, kind: str
) -> list[dict[str, Any]]:
    groups: dict[str, list[dict[str, Any]]] = defaultdict(list)
    for document in documents:
        value = document.get(field)
        if value:
            groups[str(value)].append(document)
    results: list[dict[str, Any]] = []
    for digest, members in groups.items():
        if len(members) < 2:
            continue
        newest = max(members, key=lambda item: item["modifiedAt"])
        results.append(
            {
                "kind": kind,
                "hash": digest,
                "documents": [member["id"] for member in members],
                "suggestedPrimary": newest["id"],
                "reason": (
                    "File bytes are identical."
                    if kind == "exact"
                    else "Normalized extracted text is identical."
                ),
            }
        )
    return results


def cosine(left: Iterable[float], right: Iterable[float]) -> float:
    left_values = list(map(float, left))
    right_values = list(map(float, right))
    numerator = sum(a * b for a, b in zip(left_values, right_values))
    left_norm = math.sqrt(sum(value * value for value in left_values))
    right_norm = math.sqrt(sum(value * value for value in right_values))
    return numerator / (left_norm * right_norm) if left_norm and right_norm else 0.0


def build_embeddings(
    texts: list[str], config: dict[str, Any], warnings: list[str]
) -> list[list[float]] | None:
    if not config.get("useEmbeddings", True):
        return None
    try:
        from fastembed import TextEmbedding

        model = TextEmbedding(model_name=config["embeddingModel"])
        limit = int(config.get("maximumEmbeddingCharacters", 12000))
        return [list(vector) for vector in model.embed([text[:limit] for text in texts])]
    except Exception as exc:
        warnings.append(
            "Local embeddings unavailable; lexical analysis continued. "
            f"{type(exc).__name__}: {exc}"
        )
        return None


def bm25_scores(
    tokenized: list[list[str]], warnings: list[str]
) -> tuple[list[list[float]], bool]:
    try:
        from rank_bm25 import BM25Okapi

        model = BM25Okapi(tokenized)
        rows: list[list[float]] = []
        for index, tokens in enumerate(tokenized):
            raw = list(map(float, model.get_scores(tokens)))
            candidates = [score for position, score in enumerate(raw) if position != index]
            scale = max(candidates, default=0.0)
            rows.append([score / scale if scale > 0 else 0.0 for score in raw])
        return rows, True
    except Exception as exc:
        warnings.append(
            "BM25 unavailable; token-overlap analysis continued. "
            f"{type(exc).__name__}: {exc}"
        )
        return [[0.0 for _ in tokenized] for _ in tokenized], False


def sentence_units(text: str, pages: list[str]) -> list[dict[str, Any]]:
    sources = list(enumerate(pages, start=1)) if pages else [(None, text)]
    units: list[dict[str, Any]] = []
    for page, source in sources:
        for sentence in SENTENCE_RE.split(source):
            excerpt = " ".join(sentence.split()).strip()
            if len(excerpt) >= 25:
                units.append({"page": page, "excerpt": excerpt})
    return units


def topic_tokens(sentence: str) -> set[str]:
    return {token for token in tokenize(sentence) if token not in TOPIC_STOPWORDS}


def topic_overlap(left: str, right: str) -> tuple[int, float]:
    left_tokens = topic_tokens(left)
    right_tokens = topic_tokens(right)
    shared = left_tokens & right_tokens
    union = left_tokens | right_tokens
    return len(shared), len(shared) / len(union) if union else 0.0


def evidence_record(
    evidence_type: str,
    left: dict[str, Any],
    right: dict[str, Any],
) -> dict[str, Any]:
    return {
        "type": evidence_type,
        "leftPage": left["page"] or "",
        "leftExcerpt": left["excerpt"][:500],
        "rightPage": right["page"] or "",
        "rightExcerpt": right["excerpt"][:500],
    }


def opposing_directive_evidence(
    left_units: list[dict[str, Any]],
    right_units: list[dict[str, Any]],
    terms: list[str],
) -> dict[str, Any] | None:
    lowered_terms = tuple(term.lower() for term in terms)
    left_candidates = [
        unit
        for unit in left_units
        if any(term in unit["excerpt"].lower() for term in lowered_terms)
    ]
    right_candidates = [
        unit
        for unit in right_units
        if any(term in unit["excerpt"].lower() for term in lowered_terms)
    ]
    best: tuple[float, dict[str, Any], dict[str, Any]] | None = None
    for left in left_candidates:
        left_polarity = directive_polarity(left["excerpt"])
        if left_polarity == "neutral":
            continue
        for right in right_candidates:
            right_polarity = directive_polarity(right["excerpt"])
            if {left_polarity, right_polarity} != {"positive", "negative"}:
                continue
            shared, overlap = topic_overlap(left["excerpt"], right["excerpt"])
            if shared >= 2 and overlap >= 0.3:
                score = shared + overlap
                if best is None or score > best[0]:
                    best = (score, left, right)
    return evidence_record("opposing_directive", best[1], best[2]) if best else None


def differing_value_evidence(
    left_units: list[dict[str, Any]],
    right_units: list[dict[str, Any]],
    pattern: re.Pattern[str],
    evidence_type: str,
) -> dict[str, Any] | None:
    left_candidates = [
        (unit, set(pattern.findall(unit["excerpt"])))
        for unit in left_units
        if pattern.search(unit["excerpt"])
    ]
    right_candidates = [
        (unit, set(pattern.findall(unit["excerpt"])))
        for unit in right_units
        if pattern.search(unit["excerpt"])
    ]
    best: tuple[float, dict[str, Any], dict[str, Any]] | None = None
    for left, left_values in left_candidates:
        for right, right_values in right_candidates:
            if left_values == right_values:
                continue
            shared, overlap = topic_overlap(left["excerpt"], right["excerpt"])
            if shared >= 3 and overlap >= 0.25:
                score = shared + overlap
                if best is None or score > best[0]:
                    best = (score, left, right)
    return evidence_record(evidence_type, best[1], best[2]) if best else None


def conflict_evidence(
    left: str,
    right: str,
    similarity: float,
    terms: list[str],
    left_pages: list[str],
    right_pages: list[str],
) -> tuple[list[str], list[dict[str, Any]]]:
    left_units = sentence_units(left, left_pages)
    right_units = sentence_units(right, right_pages)
    evidence: list[dict[str, Any]] = []
    directive = opposing_directive_evidence(left_units, right_units, terms)
    if directive:
        evidence.append(directive)
    date = differing_value_evidence(left_units, right_units, DATE_RE, "date")
    if date:
        evidence.append(date)
    if similarity >= 0.8:
        number = differing_value_evidence(left_units, right_units, NUMBER_RE, "numeric")
        if number:
            evidence.append(number)
    labels = {
        "opposing_directive": "Opposing directive language appears on the same topic.",
        "date": "Comparable passages contain different explicit dates.",
        "numeric": "Comparable passages contain different numeric values.",
    }
    return [labels[item["type"]] for item in evidence], evidence


def directive_polarity(sentence: str) -> str:
    lowered = sentence.lower()
    if any(term in lowered for term in NEGATIVE_DIRECTIVES):
        return "negative"
    if any(term in lowered for term in POSITIVE_DIRECTIVES):
        return "positive"
    return "neutral"


def directive_topic_tokens(sentence: str) -> set[str]:
    return topic_tokens(sentence)


def similarity_analysis(
    documents: list[dict[str, Any]],
    extracted: dict[str, str],
    page_texts: dict[str, list[str]],
    config: dict[str, Any],
    warnings: list[str],
) -> tuple[list[dict[str, Any]], list[dict[str, Any]], list[str]]:
    representatives: dict[str, dict[str, Any]] = {}
    for document in documents:
        if document["id"] not in extracted:
            continue
        key = document.get("normalizedSha256") or document["id"]
        current = representatives.get(key)
        if current is None or document["modifiedAt"] > current["modifiedAt"]:
            representatives[key] = document
    eligible = list(representatives.values())
    maximum_documents = int(config.get("maximumPairwiseDocuments", 500))
    if len(eligible) > maximum_documents:
        warnings.append(
            "Pairwise similarity and conflict analysis skipped: "
            f"{len(eligible)} extracted documents exceeds the configured "
            f"maximumPairwiseDocuments value of {maximum_documents}."
        )
        return [], [], ["pairwise-analysis-skipped"]
    texts = [extracted[document["id"]] for document in eligible]
    tokens = [tokenize(text) for text in texts]
    token_sets = [set(row) for row in tokens]
    bm25, bm25_available = bm25_scores(tokens, warnings)
    embeddings = build_embeddings(texts, config, warnings)
    methods = ["token-overlap"]
    if bm25_available:
        methods.append("bm25")
    if embeddings is not None:
        methods.append("local-embeddings")

    near_threshold = float(config["nearDuplicateThreshold"])
    related_threshold = float(config["relatedContentThreshold"])
    conflict_threshold = float(config["conflictCandidateThreshold"])
    terms = list(config.get("directiveTerms", []))
    similarities: list[dict[str, Any]] = []
    conflicts: list[dict[str, Any]] = []
    for left_index in range(len(eligible)):
        for right_index in range(left_index + 1, len(eligible)):
            left_set, right_set = token_sets[left_index], token_sets[right_index]
            union = left_set | right_set
            jaccard = len(left_set & right_set) / len(union) if union else 0.0
            bm25_pair = (
                bm25[left_index][right_index] + bm25[right_index][left_index]
            ) / 2
            lexical = (jaccard + bm25_pair) / 2 if bm25_available else jaccard
            semantic = (
                cosine(embeddings[left_index], embeddings[right_index])
                if embeddings is not None
                else lexical
            )
            combined = 0.7 * semantic + 0.3 * lexical if embeddings is not None else lexical
            if combined < related_threshold:
                continue
            kind = "near_duplicate" if combined >= near_threshold else "related_content"
            candidate = {
                "left": eligible[left_index]["id"],
                "right": eligible[right_index]["id"],
                "kind": kind,
                "score": round(combined, 4),
                "semanticScore": round(semantic, 4),
                "lexicalScore": round(lexical, 4),
            }
            similarities.append(candidate)
            if combined >= conflict_threshold:
                reasons, evidence = conflict_evidence(
                    texts[left_index],
                    texts[right_index],
                    combined,
                    terms,
                    page_texts.get(candidate["left"], []),
                    page_texts.get(candidate["right"], []),
                )
                if reasons:
                    conflicts.append(
                        {**candidate, "reasons": reasons, "evidence": evidence}
                    )
    return similarities, conflicts, methods


def stale_candidates(
    documents: list[dict[str, Any]], stale_after_days: int
) -> list[dict[str, Any]]:
    now = datetime.now(timezone.utc)
    results: list[dict[str, Any]] = []
    for document in documents:
        modified = datetime.fromisoformat(document["modifiedAt"])
        age_days = (now - modified).days
        if age_days >= stale_after_days:
            results.append(
                {
                    "document": document["id"],
                    "ageDays": age_days,
                    "thresholdDays": stale_after_days,
                }
            )
    return sorted(results, key=lambda item: item["ageDays"], reverse=True)


def comparison_evidence(
    left_id: str,
    right_id: str,
    extracted: dict[str, str],
    page_texts: dict[str, list[str]],
) -> dict[str, Any]:
    left_units = sentence_units(
        extracted.get(left_id, ""), page_texts.get(left_id, [])
    )[:400]
    right_units = sentence_units(
        extracted.get(right_id, ""), page_texts.get(right_id, [])
    )[:400]
    right_by_normalized: dict[str, dict[str, Any]] = {}
    for unit in right_units:
        normalized = normalize_text(unit["excerpt"])
        if len(normalized) >= 40:
            right_by_normalized.setdefault(normalized, unit)
    for left in left_units:
        normalized = normalize_text(left["excerpt"])
        if normalized in right_by_normalized and len(normalized) >= 40:
            return evidence_record(
                "matching_content", left, right_by_normalized[normalized]
            )
    best: tuple[float, dict[str, Any], dict[str, Any]] | None = None
    for left in left_units:
        for right in right_units:
            shared, overlap = topic_overlap(left["excerpt"], right["excerpt"])
            if shared >= 5 and overlap >= 0.4:
                score = shared + overlap
                if best is None or score > best[0]:
                    best = (score, left, right)
    return (
        evidence_record("similar_content", best[1], best[2])
        if best
        else {
            "type": "whole_document",
            "leftPage": "",
            "leftExcerpt": "The complete extracted documents matched; no concise passage was isolated.",
            "rightPage": "",
            "rightExcerpt": "The complete extracted documents matched; no concise passage was isolated.",
        }
    )


def build_backlog(
    documents: list[dict[str, Any]],
    duplicate_groups: list[dict[str, Any]],
    similarities: list[dict[str, Any]],
    conflicts: list[dict[str, Any]],
    stale: list[dict[str, Any]],
    extracted: dict[str, str],
    page_texts: dict[str, list[str]],
) -> list[dict[str, Any]]:
    by_id = {document["id"]: document for document in documents}
    backlog: list[dict[str, Any]] = []

    def add(
        priority: str,
        category: str,
        action: str,
        primary: str,
        related: str = "",
        confidence: float | None = None,
        reason: str = "",
        evidence: dict[str, Any] | None = None,
    ) -> None:
        document = by_id[primary]
        related_document = by_id[related] if related else {}
        evidence = evidence or {}
        backlog.append(
            {
                "priority": priority,
                "category": category,
                "recommendedAction": action,
                "primaryDocument": document["name"],
                "primaryPath": document["path"],
                "primaryPage": evidence.get("leftPage", ""),
                "primaryExcerpt": evidence.get("leftExcerpt", ""),
                "relatedDocument": related_document.get("name", ""),
                "relatedPath": related_document.get("path", ""),
                "relatedPage": evidence.get("rightPage", ""),
                "relatedExcerpt": evidence.get("rightExcerpt", ""),
                "confidence": confidence,
                "reason": reason,
                "sourceUrl": document.get("sourceUrl", ""),
                "owner": document.get("owner", ""),
            }
        )

    for conflict in conflicts:
        add(
            "Critical",
            "Potential conflict",
            "Compare the cited directives and confirm the authoritative scope.",
            conflict["left"],
            conflict["right"],
            conflict["score"],
            " ".join(conflict["reasons"]),
            conflict.get("evidence", [{}])[0],
        )
    conflict_pairs = {
        tuple(sorted((conflict["left"], conflict["right"]))) for conflict in conflicts
    }
    for document in documents:
        if document["extractionStatus"] != "ok":
            add(
                "High",
                "Extraction gap",
                "Resolve extraction or OCR before treating the review as complete.",
                document["id"],
                reason=document["error"] or document["extractionStatus"],
            )
    for group in duplicate_groups:
        for document_id in group["documents"]:
            if document_id == group["suggestedPrimary"]:
                continue
            add(
                "High",
                "Duplicate",
                "Confirm the primary copy, then review this copy for merge or archive.",
                document_id,
                group["suggestedPrimary"],
                1.0,
                group["reason"],
                comparison_evidence(
                    document_id,
                    group["suggestedPrimary"],
                    extracted,
                    page_texts,
                ),
            )
    for candidate in similarities:
        pair = tuple(sorted((candidate["left"], candidate["right"])))
        if pair in conflict_pairs:
            continue
        if candidate["kind"] == "near_duplicate":
            add(
                "Medium",
                "Near duplicate",
                "Compare ownership and scope; merge or document why both are needed.",
                candidate["left"],
                candidate["right"],
                candidate["score"],
                "High content similarity.",
                comparison_evidence(
                    candidate["left"],
                    candidate["right"],
                    extracted,
                    page_texts,
                ),
            )
        elif candidate["kind"] == "related_content":
            add(
                "Low",
                "Related content",
                "Review whether a cross-link or shared authoritative source would help.",
                candidate["left"],
                candidate["right"],
                candidate["score"],
                "The documents cover strongly related content.",
            )
    for item in stale:
        add(
            "Medium",
            "Stale candidate",
            "Ask the content owner to confirm, refresh, reapprove, or archive.",
            item["document"],
            confidence=None,
            reason=f"{item['ageDays']} days since last known modification.",
        )
    return sort_backlog(backlog)


def sort_backlog(backlog: list[dict[str, Any]]) -> list[dict[str, Any]]:
    return sorted(
        backlog,
        key=lambda item: (
            PRIORITY_ORDER.get(item.get("priority", ""), 99),
            CATEGORY_ORDER.get(item.get("category", ""), 99),
        ),
    )


def display_documents(documents: list[dict[str, Any]]) -> list[dict[str, Any]]:
    return [
        {key: value for key, value in document.items() if key != "id"}
        for document in documents
    ]


def display_duplicate_groups(
    groups: list[dict[str, Any]],
    documents: list[dict[str, Any]],
    extracted: dict[str, str],
    page_texts: dict[str, list[str]],
) -> list[dict[str, Any]]:
    by_id = {document["id"]: document for document in documents}
    displayed: list[dict[str, Any]] = []
    for group in groups:
        primary_id = group["suggestedPrimary"]
        comparison_id = next(
            document_id
            for document_id in group["documents"]
            if document_id != primary_id
        )
        evidence = comparison_evidence(
            comparison_id, primary_id, extracted, page_texts
        )
        displayed.append(
            {
                "kind": group["kind"],
                "hash": group["hash"],
                "documents": [by_id[item]["name"] for item in group["documents"]],
                "documentPaths": [
                    by_id[item]["path"] for item in group["documents"]
                ],
                "suggestedPrimary": by_id[primary_id]["name"],
                "suggestedPrimaryPath": by_id[primary_id]["path"],
                "evidencePage": evidence.get("rightPage", ""),
                "evidenceExcerpt": evidence.get("rightExcerpt", ""),
                "reason": group["reason"],
            }
        )
    return displayed


def display_pair_records(
    records: list[dict[str, Any]],
    documents: list[dict[str, Any]],
    extracted: dict[str, str],
    page_texts: dict[str, list[str]],
) -> list[dict[str, Any]]:
    by_id = {document["id"]: document for document in documents}
    displayed: list[dict[str, Any]] = []
    for record in records:
        left_id, right_id = record["left"], record["right"]
        evidence_items = record.get("evidence", [])
        evidence = evidence_items[0] if evidence_items else {}
        if not evidence and record.get("kind") == "near_duplicate":
            evidence = comparison_evidence(
                left_id, right_id, extracted, page_texts
            )
        output = {
            "leftDocument": by_id[left_id]["name"],
            "leftPath": by_id[left_id]["path"],
            "leftPage": evidence.get("leftPage", ""),
            "leftExcerpt": evidence.get("leftExcerpt", ""),
            "rightDocument": by_id[right_id]["name"],
            "rightPath": by_id[right_id]["path"],
            "rightPage": evidence.get("rightPage", ""),
            "rightExcerpt": evidence.get("rightExcerpt", ""),
            "kind": record["kind"],
            "score": record["score"],
            "semanticScore": record["semanticScore"],
            "lexicalScore": record["lexicalScore"],
        }
        if record.get("reasons"):
            output["reasons"] = record["reasons"]
        if evidence_items:
            output["evidence"] = evidence_items
        displayed.append(output)
    return displayed


def display_stale(
    stale: list[dict[str, Any]], documents: list[dict[str, Any]]
) -> list[dict[str, Any]]:
    by_id = {document["id"]: document for document in documents}
    return [
        {
            "document": by_id[item["document"]]["name"],
            "path": by_id[item["document"]]["path"],
            "ageDays": item["ageDays"],
            "thresholdDays": item["thresholdDays"],
        }
        for item in stale
    ]


def write_workbook(
    output: Path,
    documents: list[dict[str, Any]],
    backlog: list[dict[str, Any]],
    summary: dict[str, Any],
    config: dict[str, Any],
    analysis_methods: list[str],
    warnings: list[str],
    metadata_supplied: bool,
    corpus_scope: str | None,
    content_scope: str | None,
    source_batches: list[str],
) -> None:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill
    from openpyxl.utils import get_column_letter

    workbook = Workbook()
    workbook.remove(workbook.active)
    sorted_backlog = sort_backlog(backlog)
    sheets = {
        "Review Backlog": records_to_rows(sorted_backlog, BACKLOG_COLUMNS),
        "Summary": [["Metric", "Value"], *[[key, value] for key, value in summary.items()]],
        "Document Inventory": records_to_rows(documents, DOCUMENT_COLUMNS),
        "Curation Settings": curation_settings_rows(
            config,
            analysis_methods,
            warnings,
            metadata_supplied,
            corpus_scope,
            content_scope,
            source_batches,
        ),
    }
    for name, rows in sheets.items():
        sheet = workbook.create_sheet(name)
        for row in rows:
            sheet.append([serialize_cell(value) for value in row])
        if sheet.max_row:
            for cell in sheet[1]:
                cell.font = Font(name="Arial", size=10, bold=True, color="FFFFFF")
                cell.fill = PatternFill("solid", fgColor="1F4E78")
                cell.alignment = Alignment(wrap_text=True)
            sheet.freeze_panes = "A2"
            sheet.auto_filter.ref = sheet.dimensions
            for column in range(1, sheet.max_column + 1):
                width = max(
                    len(str(sheet.cell(row, column).value or ""))
                    for row in range(1, min(sheet.max_row, 200) + 1)
                )
                sheet.column_dimensions[get_column_letter(column)].width = min(
                    max(width + 2, 12), 55
                )
            for row in sheet.iter_rows(min_row=2):
                for cell in row:
                    cell.font = Font(name="Arial", size=10)
                    cell.alignment = Alignment(vertical="top", wrap_text=True)
            if name == "Review Backlog":
                confidence_column = next(
                    (
                        cell.column
                        for cell in sheet[1]
                        if str(cell.value).strip().lower() == "confidence"
                    ),
                    None,
                )
                if confidence_column:
                    for row in range(2, sheet.max_row + 1):
                        sheet.cell(row, confidence_column).number_format = "0.0%"
    workbook.save(output)


def curation_settings_rows(
    config: dict[str, Any],
    analysis_methods: list[str],
    warnings: list[str],
    metadata_supplied: bool,
    corpus_scope: str | None,
    content_scope: str | None,
    source_batches: list[str],
) -> list[list[Any]]:
    scope_interpretation = {
        "whole-library": (
            "The user identified the upload as the whole intended library. "
            "Completion is still stated as Complete for uploaded corpus because "
            "SharePoint coverage was not independently verified."
        ),
        "subset": (
            "The user identified the upload as a subset of the library. Findings "
            "and completion apply only to the uploaded corpus."
        ),
    }.get(
        corpus_scope,
        "Library coverage was not stated. Findings and completion apply only to "
        "the uploaded corpus.",
    )
    content_scope_interpretation = {
        "current-only": (
            "Current content only; drafts, archives, and historical versions were "
            "excluded at the user's direction."
        ),
        "include-drafts-and-history": (
            "Drafts, archives, and historical versions were included at the "
            "user's direction."
        ),
    }.get(
        content_scope,
        "Treatment of drafts, archives, and historical versions was not stated.",
    )
    freshness_basis = (
        "Metadata manifest supplied; freshness uses manifest modified dates when "
        "available."
        if metadata_supplied
        else "No metadata manifest supplied. Downloaded or staged timestamps can "
        "replace original SharePoint modified dates, so freshness findings are "
        "provisional."
    )
    rows = [
        ["Area", "Setting / interpretation"],
        [
            "Source uploads",
            ", ".join(source_batches) if source_batches else "Batch manifest not supplied.",
        ],
        ["Scope", scope_interpretation],
        ["Content versions", content_scope_interpretation],
        [
            "Freshness",
            f"Stale threshold: {config['staleAfterDays']} days. {freshness_basis}",
        ],
        [
            "Similarity",
            "Methods: "
            + (", ".join(analysis_methods) if analysis_methods else "none")
            + ". Similarity is a candidate-generation signal, not proof.",
        ],
        [
            "Thresholds",
            f"Near duplicate: {config['nearDuplicateThreshold']}; related content: "
            f"{config['relatedContentThreshold']}; potential conflict: "
            f"{config['conflictCandidateThreshold']}.",
        ],
        [
            "Pairwise analysis",
            f"Maximum documents: {config['maximumPairwiseDocuments']}.",
        ],
        [
            "Curation",
            "No source file was changed, deleted, moved, renamed, archived, or "
            "published. Every action is a human-review recommendation.",
        ],
    ]
    rows.extend(
        [["Warning", warning] for warning in warnings]
        or [["Warnings", "No deterministic analysis warnings."]]
    )
    return rows


def serialize_cell(value: Any) -> Any:
    if isinstance(value, (dict, list)):
        value = json.dumps(value, ensure_ascii=False)
    if value is None:
        return ""
    if isinstance(value, str) and value.startswith(("=", "+", "-", "@")):
        return f"'{value}"
    return value


def records_to_rows(
    records: list[dict[str, Any]], columns: Iterable[str]
) -> list[list[Any]]:
    headers = list(columns)
    return [headers, *[[record.get(header, "") for header in headers] for record in records]]


def write_html_report(
    output: Path,
    summary: dict[str, Any],
    backlog: list[dict[str, Any]],
    warnings: list[str],
    generated_at: datetime,
    source_batches: list[str],
) -> None:
    sorted_backlog = sort_backlog(backlog)
    created_label = generated_at.strftime("%Y-%m-%d %H:%M:%S UTC")
    source_label = ", ".join(source_batches) if source_batches else "Batch manifest not supplied"
    cards = "".join(
        f"<div class='card'><strong>{html.escape(str(key))}</strong>"
        f"<span>{html.escape(str(value))}</span></div>"
        for key, value in summary.items()
    )
    rows = "".join(
        "<tr>"
        + "".join(
            f"<td>{html.escape(str(item.get(field, '') or ''))}</td>"
            for field, _label in REPORT_COLUMNS
        )
        + "</tr>"
        for item in sorted_backlog[:100]
    )
    warnings_html = "".join(f"<li>{html.escape(item)}</li>" for item in warnings)
    truncation = (
        f"<p><em>Showing the top 100 of {len(sorted_backlog)} backlog items. "
        "See the Excel workbook for the complete list.</em></p>"
        if len(sorted_backlog) > 100
        else ""
    )
    output.write_text(
        f"""<!doctype html>
<html lang="en"><head><meta charset="utf-8">
<title>Knowledge Corpus Curation Report - {html.escape(created_label)}</title>
<style>
body{{font-family:Segoe UI,Arial,sans-serif;margin:40px;color:#172033;background:#f5f7fb}}
h1{{color:#143d66}} .cards{{display:flex;flex-wrap:wrap;gap:12px}}
.card{{background:white;border-radius:8px;padding:14px;min-width:150px;box-shadow:0 1px 4px #ccd}}
.card strong,.card span{{display:block}} .card span{{font-size:24px;margin-top:6px}}
table{{border-collapse:collapse;width:100%;background:white;margin-top:20px}}
th,td{{border:1px solid #d8deea;padding:8px;vertical-align:top;text-align:left}}
th{{background:#1f4e78;color:white}} li{{margin:6px 0}}
</style></head><body>
<h1>Knowledge Corpus Curation Report - {html.escape(created_label)}</h1>
<p><strong>Source ZIP file(s):</strong> {html.escape(source_label)}</p>
<div class="cards">{cards}</div>
<h2>Priority backlog</h2>
{truncation}
<table><thead><tr>{''.join(f'<th>{html.escape(label)}</th>' for _field, label in REPORT_COLUMNS)}</tr></thead>
<tbody>{rows}</tbody></table>
<h2>Warnings</h2><ul>{warnings_html or '<li>None</li>'}</ul>
</body></html>""",
        encoding="utf-8",
    )


def main() -> int:
    args = parse_args()
    if not args.input.is_dir():
        raise SystemExit(f"Input directory does not exist: {args.input}")
    args.output.mkdir(parents=True, exist_ok=True)
    config = load_json(args.config, {})
    defaults = {
        "staleAfterDays": 365,
        "nearDuplicateThreshold": 0.86,
        "relatedContentThreshold": 0.68,
        "conflictCandidateThreshold": 0.72,
        "minimumExtractedCharacters": 80,
        "maximumEmbeddingCharacters": 12000,
        "maximumPairwiseDocuments": 500,
        "useEmbeddings": True,
        "embeddingModel": "BAAI/bge-small-en-v1.5",
        "directiveTerms": [
            "must",
            "must not",
            "should",
            "should not",
            "required",
            "prohibited",
            "never",
            "always",
            "optional",
        ],
    }
    defaults.update(config)
    config = defaults
    if args.stale_after_days is not None:
        config["staleAfterDays"] = args.stale_after_days
    if args.no_embeddings:
        config["useEmbeddings"] = False

    metadata = load_metadata(args.metadata)
    source_batches = load_batch_sources(args.batch_manifest)
    documents, extracted, page_texts, warnings = inventory(
        args.input,
        metadata,
        int(config["minimumExtractedCharacters"]),
        args.ocr,
    )
    duplicates = exact_duplicate_groups(documents, "sha256", "exact")
    normalized_groups = exact_duplicate_groups(
        [doc for doc in documents if doc["extractionStatus"] == "ok"],
        "normalizedSha256",
        "normalized_text",
    )
    known_exact_members = {
        document_id for group in duplicates for document_id in group["documents"]
    }
    normalized_groups = [
        group
        for group in normalized_groups
        if not set(group["documents"]).issubset(known_exact_members)
    ]
    duplicate_groups = duplicates + normalized_groups
    similarities, conflicts, methods = similarity_analysis(
        documents, extracted, page_texts, config, warnings
    )
    stale = stale_candidates(documents, int(config["staleAfterDays"]))
    backlog = sort_backlog(build_backlog(
        documents,
        duplicate_groups,
        similarities,
        conflicts,
        stale,
        extracted,
        page_texts,
    ))
    output_documents = display_documents(documents)
    output_duplicate_groups = display_duplicate_groups(
        duplicate_groups, documents, extracted, page_texts
    )
    output_similarities = display_pair_records(
        similarities, documents, extracted, page_texts
    )
    output_conflicts = display_pair_records(
        conflicts, documents, extracted, page_texts
    )
    output_stale = display_stale(stale, documents)
    status_counts = Counter(document["extractionStatus"] for document in documents)
    summary = {
        "filesScanned": len(documents),
        "filesExtracted": status_counts.get("ok", 0),
        "extractionGaps": len(documents) - status_counts.get("ok", 0),
        "duplicateGroups": len(duplicate_groups),
        "nearDuplicatePairs": sum(
            candidate["kind"] == "near_duplicate" for candidate in similarities
        ),
        "relatedContentPairs": sum(
            candidate["kind"] == "related_content" for candidate in similarities
        ),
        "conflictCandidates": len(conflicts),
        "staleCandidates": len(stale),
        "backlogItems": len(backlog),
    }
    generated_at = datetime.now(timezone.utc)
    creation_stamp = generated_at.strftime("%Y-%m-%d-%H%M%SZ")
    result = {
        "generatedAt": generated_at.isoformat(),
        "inputRoot": str(args.input.resolve()),
        "sourceBatches": source_batches,
        "analysisMode": methods,
        "warnings": warnings,
        "summary": summary,
        "documents": output_documents,
        "duplicateGroups": output_duplicate_groups,
        "similarityCandidates": output_similarities,
        "conflictCandidates": output_conflicts,
        "staleCandidates": output_stale,
        "backlog": backlog,
    }
    json_path = args.output / "curation-results.json"
    workbook_path = (
        args.output / f"knowledge-corpus-curation-backlog-{creation_stamp}.xlsx"
    )
    html_path = (
        args.output / f"knowledge-corpus-curation-report-{creation_stamp}.html"
    )
    json_path.write_text(json.dumps(result, indent=2, ensure_ascii=False), encoding="utf-8")
    write_workbook(
        workbook_path,
        output_documents,
        backlog,
        summary,
        config,
        methods,
        warnings,
        args.metadata is not None,
        args.corpus_scope,
        args.content_scope,
        source_batches,
    )
    write_html_report(
        html_path, summary, backlog, warnings, generated_at, source_batches
    )
    print(
        json.dumps(
            {
                "json": str(json_path),
                "workbook": str(workbook_path),
                "html": str(html_path),
                "summary": summary,
                "analysisMode": methods,
                "warnings": warnings,
            },
            ensure_ascii=False,
        )
    )
    return 0


if __name__ == "__main__":
    sys.exit(main())
