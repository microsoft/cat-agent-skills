#!/usr/bin/env python3
"""Universal Document Converter — offline document format conversion.

Usage:
    python convert.py INPUT --to FORMAT [-o OUTPUT]
    python convert.py --batch DIR --to FORMAT [--out-dir DIR]

Formats: md, html, pdf, docx, pptx, txt (inputs also: xlsx, csv).

Every pipeline uses only libraries available in the agent sandbox
(markitdown, mammoth, markdownify, reportlab, python-docx, python-pptx,
pdfplumber, beautifulsoup4, magika) — no network access required.
"""

from __future__ import annotations

import argparse
import os
import sys

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

import blocks as blocks_mod
import render

INPUT_FORMATS = {"md", "html", "docx", "pptx", "xlsx", "csv", "pdf", "txt"}
OUTPUT_FORMATS = {"md", "html", "pdf", "docx", "pptx", "txt"}

# Explicit support matrix: input format -> allowed output formats.
# pdf/pptx/xlsx/csv -> pptx is deliberately unsupported: slides built from
# extraction artifacts are too lossy to present as a finished conversion.
SUPPORTED: dict[str, set[str]] = {
    "md":   {"html", "pdf", "docx", "pptx", "txt"},
    "html": {"md", "pdf", "docx", "pptx", "txt"},
    "docx": {"md", "html", "pdf", "pptx", "txt"},
    "pptx": {"md", "html", "txt"},
    "xlsx": {"md", "html", "pdf", "docx", "txt"},
    "csv":  {"md", "html", "pdf", "docx", "txt"},
    "pdf":  {"md", "html", "docx", "txt"},
    "txt":  {"md", "html", "pdf", "docx"},
}

_EXT_TO_FORMAT = {
    ".md": "md", ".markdown": "md", ".html": "html", ".htm": "html",
    ".docx": "docx", ".pptx": "pptx", ".xlsx": "xlsx", ".csv": "csv",
    ".pdf": "pdf", ".txt": "txt", ".text": "txt",
}


def support_matrix() -> str:
    lines = ["Supported conversions (input -> outputs):"]
    for fmt in sorted(SUPPORTED):
        lines.append(f"  {fmt:<5} -> {', '.join(sorted(SUPPORTED[fmt]))}")
    return "\n".join(lines)


class ConversionError(Exception):
    pass


# --------------------------------------------------------------------------
# Input-format detection
# --------------------------------------------------------------------------

def _sniff_content(path: str) -> str | None:
    """Best-effort content-based format detection. Returns a format or None."""
    try:  # magika ships its model inside the wheel — fully offline
        from magika import Magika
        label = str(Magika().identify_path(path).output.label)
        return {
            "markdown": "md", "html": "html", "pdf": "pdf", "csv": "csv",
            "docx": "docx", "pptx": "pptx", "xlsx": "xlsx", "txt": "txt",
        }.get(label)
    except Exception:
        pass
    try:  # magic-bytes fallback
        head = open(path, "rb").read(512)
    except OSError:
        return None
    if head.startswith(b"%PDF"):
        return "pdf"
    lowered = head.lstrip().lower()
    if lowered.startswith((b"<!doctype", b"<html")):
        return "html"
    if head.startswith(b"PK"):
        return None  # some zip-based Office format; trust the extension
    return None


def detect_input_format(path: str) -> str:
    ext_fmt = _EXT_TO_FORMAT.get(os.path.splitext(path)[1].lower())
    content_fmt = _sniff_content(path)
    if ext_fmt and content_fmt and ext_fmt != content_fmt:
        render.warn(
            f"'{os.path.basename(path)}' has a .{ext_fmt} extension but its "
            f"content looks like {content_fmt} — converting it as {content_fmt}."
        )
        return content_fmt
    fmt = ext_fmt or content_fmt
    if not fmt:
        raise ConversionError(
            f"Could not determine the format of '{path}'.\n\n{support_matrix()}"
        )
    return fmt


# --------------------------------------------------------------------------
# X -> Markdown (the universal intermediate)
# --------------------------------------------------------------------------

def _csv_to_markdown(path: str) -> str:
    import csv

    with open(path, newline="", encoding="utf-8-sig") as fh:
        rows = [row for row in csv.reader(fh) if row]
    if not rows:
        return ""
    table = blocks_mod.Table(rows)
    return render.to_markdown([table])


def to_markdown_text(path: str, fmt: str) -> str:
    if fmt == "md":
        return open(path, encoding="utf-8").read()
    if fmt == "txt":
        return open(path, encoding="utf-8").read()
    if fmt == "html":
        try:
            from markdownify import markdownify
            return markdownify(open(path, encoding="utf-8").read(), heading_style="ATX")
        except ImportError:
            return render.to_markdown(
                blocks_mod.parse_html(open(path, encoding="utf-8").read()))
    if fmt == "csv":
        return _csv_to_markdown(path)
    # docx / pptx / xlsx / pdf — markitdown handles them all
    try:
        from markitdown import MarkItDown
        result = MarkItDown(enable_plugins=False).convert(path)
        return result.text_content
    except Exception as exc:  # missing optional deps, unreadable file, …
        render.warn(f"markitdown could not convert {fmt} "
                    f"({exc.__class__.__name__}); using the {fmt} fallback")
    if fmt == "docx":
        import mammoth
        with open(path, "rb") as fh:
            html = mammoth.convert_to_html(fh).value
        from markdownify import markdownify
        return markdownify(html, heading_style="ATX")
    if fmt == "pdf":
        return _pdf_to_text(path)
    if fmt == "xlsx":
        return _xlsx_to_markdown(path)
    if fmt == "pptx":
        return _pptx_to_markdown(path)
    raise ConversionError(f"No working converter for {fmt}.")


def _xlsx_to_markdown(path: str) -> str:
    import openpyxl

    wb = openpyxl.load_workbook(path, data_only=True)
    doc: list = []
    for ws in wb.worksheets:
        rows = [["" if c is None else str(c) for c in row]
                for row in ws.iter_rows(values_only=True)]
        rows = [r for r in rows if any(c.strip() for c in r)]
        if not rows:
            continue
        if len(wb.worksheets) > 1:
            doc.append(blocks_mod.Heading(2, [blocks_mod.Run(ws.title)]))
        doc.append(blocks_mod.Table(rows))
    return render.to_markdown(doc)


def _pptx_to_markdown(path: str) -> str:
    from pptx import Presentation

    prs = Presentation(path)
    doc: list = []
    for slide in prs.slides:
        title = slide.shapes.title
        if title is not None and title.text.strip():
            doc.append(blocks_mod.Heading(1, [blocks_mod.Run(title.text.strip())]))
        for shape in slide.shapes:
            if shape is title or not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if text:
                    doc.append(blocks_mod.Paragraph([blocks_mod.Run(text)]))
    return render.to_markdown(doc)


def _pdf_to_text(path: str) -> str:
    import pdfplumber

    pages = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            pages.append(page.extract_text() or "")
    return "\n\n".join(pages)


# --------------------------------------------------------------------------
# Conversion driver
# --------------------------------------------------------------------------

def _title_for(path: str) -> str:
    return os.path.splitext(os.path.basename(path))[0].replace("-", " ").replace("_", " ").title()


def convert(in_path: str, target: str, out_path: str | None = None) -> str:
    if not os.path.isfile(in_path):
        raise ConversionError(f"Input file not found: {in_path}")
    target = target.lstrip(".").lower()
    if target not in OUTPUT_FORMATS:
        raise ConversionError(
            f"Unknown output format '{target}'.\n\n{support_matrix()}")

    fmt = detect_input_format(in_path)
    if fmt == target:
        raise ConversionError(f"'{in_path}' is already {fmt}; nothing to convert.")
    if target not in SUPPORTED.get(fmt, set()):
        raise ConversionError(
            f"Converting {fmt} -> {target} is not supported.\n\n{support_matrix()}")

    out_path = out_path or os.path.splitext(in_path)[0] + "." + target
    title = _title_for(in_path)

    # Direct single-library routes first (higher fidelity than via blocks).
    if target == "txt" and fmt == "pdf":
        text = _pdf_to_text(in_path)
        open(out_path, "w", encoding="utf-8").write(text)
        return out_path
    if target == "html" and fmt == "docx":
        try:
            import mammoth
            with open(in_path, "rb") as fh:
                body = mammoth.convert_to_html(fh).value
            open(out_path, "w", encoding="utf-8").write(
                "<!DOCTYPE html>\n<html>\n<head>\n<meta charset=\"utf-8\">\n"
                f"<title>{title}</title>\n</head>\n<body>\n{body}\n</body>\n</html>\n")
            return out_path
        except ImportError:
            pass  # fall through to the block pipeline

    # Everything else goes via Markdown/HTML -> block model -> renderer.
    if target == "md":
        open(out_path, "w", encoding="utf-8").write(to_markdown_text(in_path, fmt))
        return out_path

    if fmt == "html":
        doc_blocks = blocks_mod.parse_html(open(in_path, encoding="utf-8").read())
    else:
        doc_blocks = blocks_mod.parse_markdown(to_markdown_text(in_path, fmt))
    if not doc_blocks:
        raise ConversionError(f"No convertible content found in '{in_path}'.")

    if target == "html":
        open(out_path, "w", encoding="utf-8").write(render.to_html(doc_blocks, title))
    elif target == "pdf":
        render.to_pdf(doc_blocks, out_path, title)
    elif target == "docx":
        render.to_docx(doc_blocks, out_path)
    elif target == "pptx":
        render.to_pptx(doc_blocks, out_path, default_title=title)
    elif target == "txt":
        open(out_path, "w", encoding="utf-8").write(render.to_text(doc_blocks))
    return out_path


def convert_batch(in_dir: str, target: str, out_dir: str | None) -> int:
    if not os.path.isdir(in_dir):
        raise ConversionError(f"Not a directory: {in_dir}")
    out_dir = out_dir or in_dir
    os.makedirs(out_dir, exist_ok=True)

    results: list[tuple[str, str, str]] = []  # (file, status, detail)
    for name in sorted(os.listdir(in_dir)):
        path = os.path.join(in_dir, name)
        if not os.path.isfile(path):
            continue
        ext_fmt = _EXT_TO_FORMAT.get(os.path.splitext(name)[1].lower())
        if ext_fmt is None:
            results.append((name, "skipped", "unrecognized extension"))
            continue
        if ext_fmt == target:
            results.append((name, "skipped", f"already {target}"))
            continue
        if target not in SUPPORTED.get(ext_fmt, set()):
            results.append((name, "skipped", f"{ext_fmt} -> {target} unsupported"))
            continue
        stem = os.path.splitext(name)[0]
        out_path = os.path.join(out_dir, stem + "." + target)
        if any(r[2] == out_path for r in results):  # e.g. a.csv and a.docx -> a.md
            out_path = os.path.join(out_dir, f"{stem}.{ext_fmt}.{target}")
        try:
            convert(path, target, out_path)
            results.append((name, "converted", out_path))
        except Exception as exc:  # keep going; report at the end
            results.append((name, "failed", str(exc).splitlines()[0]))

    width = max((len(r[0]) for r in results), default=4)
    print(f"{'FILE':<{width}}  {'STATUS':<9}  DETAIL")
    for name, status, detail in results:
        print(f"{name:<{width}}  {status:<9}  {detail}")
    converted = sum(1 for r in results if r[1] == "converted")
    failed = sum(1 for r in results if r[1] == "failed")
    print(f"\n{converted} converted, {failed} failed, "
          f"{len(results) - converted - failed} skipped.")
    return 1 if failed else 0


def main(argv: list[str] | None = None) -> int:
    parser = argparse.ArgumentParser(
        description="Convert documents between formats, fully offline.",
        epilog=support_matrix(),
        formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument("input", nargs="?", help="input file")
    parser.add_argument("--to", required=True, dest="target",
                        help="output format: " + ", ".join(sorted(OUTPUT_FORMATS)))
    parser.add_argument("-o", "--output", help="output file path")
    parser.add_argument("--batch", metavar="DIR",
                        help="convert every supported file in DIR")
    parser.add_argument("--out-dir", help="output directory for --batch")
    args = parser.parse_args(argv)

    try:
        if args.batch:
            return convert_batch(args.batch, args.target.lstrip(".").lower(), args.out_dir)
        if not args.input:
            parser.error("provide an input file or --batch DIR")
        out_path = convert(args.input, args.target, args.output)
        print(out_path)
        return 0
    except ConversionError as exc:
        print(f"error: {exc}", file=sys.stderr)
        return 2


if __name__ == "__main__":
    raise SystemExit(main())
