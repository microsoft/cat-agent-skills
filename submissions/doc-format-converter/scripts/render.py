"""Renderers for the Universal Document Converter block model.

Writes a list of blocks (see blocks.py) out as PDF (reportlab, with CJK font
registration), DOCX (python-docx), PPTX (python-pptx), standalone HTML,
Markdown, or plain text.
"""

from __future__ import annotations

import glob
import html as html_mod
import os
import re
import sys

from blocks import Block, CodeBlock, Heading, Hr, ListBlock, Paragraph, Run, Table, plain_text

_CJK_RE = re.compile(
    "["
    "\u1100-\u11ff"  # Hangul Jamo
    "\u3000-\u30ff"  # CJK punctuation, Hiragana, Katakana
    "\u3400-\u4dbf"  # CJK ideographs extension A
    "\u4e00-\u9fff"  # CJK unified ideographs
    "\uac00-\ud7af"  # Hangul syllables
    "\uf900-\ufaff"  # CJK compatibility ideographs
    "\uff00-\uffef"  # full-width forms
    "]"
)


def warn(msg: str) -> None:
    print(f"[doc-format-converter] {msg}", file=sys.stderr)


# --------------------------------------------------------------------------
# Markdown / text / HTML output
# --------------------------------------------------------------------------

def _runs_to_md(runs: list[Run]) -> str:
    parts = []
    for r in runs:
        t = r.text
        if not t:
            continue
        if r.code:
            t = f"`{t}`"
        else:
            if r.bold:
                t = f"**{t}**"
            if r.italic:
                t = f"*{t}*"
        parts.append(t)
    return "".join(parts)


def to_markdown(blocks: list[Block]) -> str:
    out: list[str] = []
    for b in blocks:
        if isinstance(b, Heading):
            out.append(f"{'#' * b.level} {_runs_to_md(b.runs)}")
        elif isinstance(b, Paragraph):
            out.append(_runs_to_md(b.runs))
        elif isinstance(b, ListBlock):
            lines = []
            for n, item in enumerate(b.items, 1):
                prefix = f"{n}." if b.ordered else "-"
                lines.append(f"{prefix} {_runs_to_md(item)}")
            out.append("\n".join(lines))
        elif isinstance(b, Table) and b.rows:
            width = max(len(r) for r in b.rows)
            rows = [r + [""] * (width - len(r)) for r in b.rows]
            lines = ["| " + " | ".join(rows[0]) + " |",
                     "| " + " | ".join(["---"] * width) + " |"]
            lines += ["| " + " | ".join(r) + " |" for r in rows[1:]]
            out.append("\n".join(lines))
        elif isinstance(b, CodeBlock):
            out.append(f"```{b.lang}\n{b.text}\n```")
        elif isinstance(b, Hr):
            out.append("---")
    return "\n\n".join(out) + "\n"


def to_text(blocks: list[Block]) -> str:
    out: list[str] = []
    for b in blocks:
        if isinstance(b, (Heading, Paragraph)):
            out.append(plain_text(b.runs))
        elif isinstance(b, ListBlock):
            out.append("\n".join(
                f"{f'{n}.' if b.ordered else '-'} {plain_text(item)}"
                for n, item in enumerate(b.items, 1)))
        elif isinstance(b, Table):
            out.append("\n".join("  ".join(r) for r in b.rows))
        elif isinstance(b, CodeBlock):
            out.append(b.text)
    return "\n\n".join(out) + "\n"


def _runs_to_html(runs: list[Run]) -> str:
    parts = []
    for r in runs:
        t = html_mod.escape(r.text)
        if r.code:
            t = f"<code>{t}</code>"
        if r.bold:
            t = f"<strong>{t}</strong>"
        if r.italic:
            t = f"<em>{t}</em>"
        parts.append(t)
    return "".join(parts)


def to_html(blocks: list[Block], title: str = "Document") -> str:
    body: list[str] = []
    for b in blocks:
        if isinstance(b, Heading):
            body.append(f"<h{b.level}>{_runs_to_html(b.runs)}</h{b.level}>")
        elif isinstance(b, Paragraph):
            body.append(f"<p>{_runs_to_html(b.runs)}</p>")
        elif isinstance(b, ListBlock):
            tag = "ol" if b.ordered else "ul"
            items = "".join(f"<li>{_runs_to_html(i)}</li>" for i in b.items)
            body.append(f"<{tag}>{items}</{tag}>")
        elif isinstance(b, Table) and b.rows:
            head = "".join(f"<th>{html_mod.escape(c)}</th>" for c in b.rows[0])
            rows = "".join(
                "<tr>" + "".join(f"<td>{html_mod.escape(c)}</td>" for c in r) + "</tr>"
                for r in b.rows[1:])
            body.append(f"<table><thead><tr>{head}</tr></thead><tbody>{rows}</tbody></table>")
        elif isinstance(b, CodeBlock):
            body.append(f"<pre><code>{html_mod.escape(b.text)}</code></pre>")
        elif isinstance(b, Hr):
            body.append("<hr>")
    return (
        "<!DOCTYPE html>\n<html>\n<head>\n<meta charset=\"utf-8\">\n"
        f"<title>{html_mod.escape(title)}</title>\n"
        "<style>body{font-family:sans-serif;max-width:48rem;margin:2rem auto;"
        "padding:0 1rem;line-height:1.5}table{border-collapse:collapse}"
        "th,td{border:1px solid #999;padding:.3rem .6rem}"
        "pre{background:#f4f4f4;padding:.8rem;overflow-x:auto}</style>\n"
        "</head>\n<body>\n" + "\n".join(body) + "\n</body>\n</html>\n"
    )


# --------------------------------------------------------------------------
# PDF (reportlab)
# --------------------------------------------------------------------------

_FONT_DIRS = [
    "/usr/share/fonts", "/usr/local/share/fonts",
    os.path.expanduser("~/.fonts"),
]


def _register_cjk_font() -> str | None:
    """Register a CJK-capable font and return its reportlab name.

    Prefers the sandbox's bundled Noto CJK TrueType fonts (noto-cjk-fonts
    package or system font dirs); falls back to reportlab's built-in
    Unicode CID fonts, which need no font files at all.
    """
    from reportlab.pdfbase import pdfmetrics
    from reportlab.pdfbase.ttfonts import TTFont

    candidates: list[str] = []
    try:  # the sandbox ships fonts in the noto-cjk-fonts package
        import noto_cjk_fonts  # type: ignore
        pkg_dir = os.path.dirname(noto_cjk_fonts.__file__)
        candidates += glob.glob(os.path.join(pkg_dir, "**", "*.[to]t[fc]"), recursive=True)
    except ImportError:
        pass
    for d in _FONT_DIRS:
        candidates += glob.glob(os.path.join(d, "**", "*Noto*CJK*.*"), recursive=True)
        candidates += glob.glob(os.path.join(d, "**", "*noto*cjk*.*"), recursive=True)

    for path in candidates:
        try:
            pdfmetrics.registerFont(TTFont("CJK", path, subfontIndex=0))
            return "CJK"
        except Exception:
            continue

    try:  # built into reportlab, works offline, viewer supplies the glyphs
        from reportlab.pdfbase.cidfonts import UnicodeCIDFont
        pdfmetrics.registerFont(UnicodeCIDFont("HeiseiKakuGo-W5"))
        return "HeiseiKakuGo-W5"
    except Exception:
        warn("no CJK-capable font available; CJK characters may render as boxes")
        return None


def _runs_to_rl(runs: list[Run], cjk_font: str | None) -> str:
    """Runs -> reportlab Paragraph inline XML markup."""
    parts = []
    for r in runs:
        t = html_mod.escape(r.text)
        if cjk_font:
            t = _CJK_RE.sub(lambda m: f'<font name="{cjk_font}">{m.group(0)}</font>', t)
        if r.code:
            t = f'<font face="Courier">{t}</font>'
        if r.bold:
            t = f"<b>{t}</b>"
        if r.italic:
            t = f"<i>{t}</i>"
        parts.append(t)
    return "".join(parts)


def to_pdf(blocks: list[Block], out_path: str, title: str = "Document") -> None:
    from reportlab.lib import colors
    from reportlab.lib.pagesizes import A4
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.lib.units import cm
    from reportlab.platypus import (HRFlowable, ListFlowable, ListItem,
                                    Paragraph as RLParagraph, Preformatted,
                                    SimpleDocTemplate, Spacer,
                                    Table as RLTable, TableStyle)

    cjk_font = None
    if any(_CJK_RE.search(plain_text(b.runs) if isinstance(b, (Heading, Paragraph))
                          else b.text if isinstance(b, CodeBlock)
                          else " ".join(" ".join(r) for r in b.rows) if isinstance(b, Table)
                          else " ".join(plain_text(i) for i in b.items) if isinstance(b, ListBlock)
                          else "")
           for b in blocks):
        cjk_font = _register_cjk_font()

    styles = getSampleStyleSheet()
    story: list = []
    for b in blocks:
        if isinstance(b, Heading):
            style = styles[f"Heading{min(b.level, 6)}"]
            story.append(RLParagraph(_runs_to_rl(b.runs, cjk_font), style))
        elif isinstance(b, Paragraph):
            story.append(RLParagraph(_runs_to_rl(b.runs, cjk_font), styles["BodyText"]))
        elif isinstance(b, ListBlock):
            items = [ListItem(RLParagraph(_runs_to_rl(i, cjk_font), styles["BodyText"]),
                              leftIndent=24)
                     for i in b.items]
            story.append(ListFlowable(items, bulletType="1" if b.ordered else "bullet"))
        elif isinstance(b, Table) and b.rows:
            width = max(len(r) for r in b.rows)
            data = [[RLParagraph(_runs_to_rl([Run(c, bold=(ri == 0))], cjk_font),
                                 styles["BodyText"])
                     for c in row + [""] * (width - len(row))]
                    for ri, row in enumerate(b.rows)]
            t = RLTable(data, hAlign="LEFT")
            t.setStyle(TableStyle([
                ("GRID", (0, 0), (-1, -1), 0.5, colors.grey),
                ("BACKGROUND", (0, 0), (-1, 0), colors.whitesmoke),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
            ]))
            story.append(Spacer(1, 0.2 * cm))
            story.append(t)
            story.append(Spacer(1, 0.2 * cm))
        elif isinstance(b, CodeBlock):
            story.append(Preformatted(b.text, styles["Code"]))
            story.append(Spacer(1, 0.2 * cm))
        elif isinstance(b, Hr):
            story.append(HRFlowable(width="100%", color=colors.grey))

    doc = SimpleDocTemplate(out_path, pagesize=A4, title=title,
                            leftMargin=2 * cm, rightMargin=2 * cm,
                            topMargin=2 * cm, bottomMargin=2 * cm)
    doc.build(story)


# --------------------------------------------------------------------------
# DOCX (python-docx)
# --------------------------------------------------------------------------

def to_docx(blocks: list[Block], out_path: str) -> None:
    import docx

    doc = docx.Document()

    def add_runs(para, runs: list[Run]) -> None:
        for r in runs:
            run = para.add_run(r.text)
            run.bold = r.bold
            run.italic = r.italic
            if r.code:
                run.font.name = "Courier New"

    for b in blocks:
        if isinstance(b, Heading):
            para = doc.add_heading(level=min(b.level, 9))
            add_runs(para, b.runs)
        elif isinstance(b, Paragraph):
            add_runs(doc.add_paragraph(), b.runs)
        elif isinstance(b, ListBlock):
            style = "List Number" if b.ordered else "List Bullet"
            for item in b.items:
                add_runs(doc.add_paragraph(style=style), item)
        elif isinstance(b, Table) and b.rows:
            width = max(len(r) for r in b.rows)
            table = doc.add_table(rows=len(b.rows), cols=width)
            table.style = "Table Grid"
            for ri, row in enumerate(b.rows):
                for ci in range(width):
                    cell = table.cell(ri, ci)
                    text = row[ci] if ci < len(row) else ""
                    run = cell.paragraphs[0].add_run(text)
                    run.bold = ri == 0
        elif isinstance(b, CodeBlock):
            para = doc.add_paragraph()
            run = para.add_run(b.text)
            run.font.name = "Courier New"

    doc.save(out_path)


# --------------------------------------------------------------------------
# PPTX (python-pptx)
# --------------------------------------------------------------------------

_MAX_BULLET_CHARS = 220
_MAX_BULLETS_PER_SLIDE = 8


def to_pptx(blocks: list[Block], out_path: str, default_title: str = "Document") -> None:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content

    slides: list[tuple[str, list[str]]] = []
    current_title: str | None = None
    current_bullets: list[str] = []

    def push() -> None:
        nonlocal current_title, current_bullets
        if current_title is not None or current_bullets:
            slides.append((current_title or default_title, current_bullets))
        current_title, current_bullets = None, []

    def bullet(text: str) -> None:
        text = text.strip()
        if not text:
            return
        if len(text) > _MAX_BULLET_CHARS:
            text = text[: _MAX_BULLET_CHARS - 1] + "…"
        current_bullets.append(text)

    for b in blocks:
        if isinstance(b, Heading) and b.level <= 2:
            push()
            current_title = plain_text(b.runs)
        elif isinstance(b, Heading):
            bullet(plain_text(b.runs))
        elif isinstance(b, Paragraph):
            bullet(plain_text(b.runs))
        elif isinstance(b, ListBlock):
            for item in b.items:
                bullet(plain_text(item))
        elif isinstance(b, Table):
            for row in b.rows:
                bullet(" | ".join(row))
        elif isinstance(b, CodeBlock):
            for line in b.text.splitlines()[:_MAX_BULLETS_PER_SLIDE]:
                bullet(line)
    push()

    if not slides:
        slides = [(default_title, [])]

    for title, bullets in slides:
        # overflow onto continuation slides rather than truncating content
        chunks = [bullets[i : i + _MAX_BULLETS_PER_SLIDE]
                  for i in range(0, len(bullets), _MAX_BULLETS_PER_SLIDE)] or [[]]
        for n, chunk in enumerate(chunks):
            slide = prs.slides.add_slide(layout)
            slide.shapes.title.text = title if n == 0 else f"{title} (cont.)"
            body = slide.placeholders[1].text_frame
            for j, text in enumerate(chunk):
                para = body.paragraphs[0] if j == 0 else body.add_paragraph()
                para.text = text
                para.font.size = Pt(18)

    prs.save(out_path)
