#!/usr/bin/env python3
"""Deterministic accessibility checks for .pptx, .docx, .html and .md files.

Findings use the same rule names as Microsoft's built-in Accessibility Checker
(Errors / Warnings / Tips), so a fix made here matches what the author sees when
they run the checker themselves in Office.

This script only reports what can be decided mechanically. Judgement calls,
such as whether alt text is *accurate*, whether a reading order makes *sense*,
or whether a caption is *correct*, are left to the agent reading these findings.

Usage:
    python scripts/a11y_check.py <file> [--json]

Requires python-pptx for .pptx and python-docx for .docx. HTML and Markdown are
handled with the standard library only.
"""

from __future__ import annotations

import argparse
import json
import os
import re
import sys
from html.parser import HTMLParser

# --- WCAG contrast -----------------------------------------------------------

# WCAG 2.1 AA: 4.5:1 for normal text, 3:1 for large text (>=18pt, or >=14pt bold).
AA_NORMAL = 4.5
AA_LARGE = 3.0


def _luminance(rgb: tuple[int, int, int]) -> float:
    def channel(v: int) -> float:
        s = v / 255
        return s / 12.92 if s <= 0.03928 else ((s + 0.055) / 1.055) ** 2.4

    r, g, b = (channel(v) for v in rgb)
    return 0.2126 * r + 0.7152 * g + 0.0722 * b


def contrast_ratio(fg: tuple[int, int, int], bg: tuple[int, int, int]) -> float:
    """Unrounded WCAG contrast ratio. Round only for display, not before a
    threshold comparison, a ratio of 4.496 must not round up to a passing
    4.50 against a 4.5 requirement."""
    a, b = _luminance(fg), _luminance(bg)
    lighter, darker = max(a, b), min(a, b)
    return (lighter + 0.05) / (darker + 0.05)


def parse_color(value: str) -> tuple[int, int, int] | None:
    """Parse #rgb, #rrggbb, or rgb(r,g,b). Returns None for anything else."""
    if not value:
        return None
    value = value.strip().lower()
    m = re.fullmatch(r"#([0-9a-f]{3})", value)
    if m:
        return tuple(int(c * 2, 16) for c in m.group(1))  # type: ignore[return-value]
    m = re.fullmatch(r"#([0-9a-f]{6})", value)
    if m:
        h = m.group(1)
        return (int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    m = re.fullmatch(r"rgba?\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)\s*(?:,[^)]*)?\)", value)
    if m:
        # CSS clamps out-of-range channel values (e.g. rgb(300,0,0)) rather
        # than rejecting them; match that instead of leaving them unbounded.
        return tuple(min(255, int(c)) for c in m.groups())  # type: ignore[return-value]
    return None


# --- findings ----------------------------------------------------------------


def finding(severity: str, rule: str, location: str, detail: str, fix: str) -> dict:
    return {
        "severity": severity,
        "rule": rule,
        "location": location,
        "detail": detail,
        "fix": fix,
    }


VAGUE_LINK_TEXT = {
    "click here", "here", "read more", "more", "link", "this", "this link",
    "learn more", "see more", "download", "click", "go", "info",
}


# --- shared XML helpers ------------------------------------------------------


def _local(tag: object) -> str:
    """Local name of an lxml tag, namespace stripped."""
    return str(tag).rsplit("}", 1)[-1]


def _find_local(element, name: str):
    """First descendant (or self) whose local name matches. Namespace-agnostic
    on purpose: the same element sits in different namespaces across pptx/docx."""
    if _local(element.tag) == name:
        return element
    for child in element.iter():
        if _local(child.tag) == name:
            return child
    return None


def _is_decorative(cnvpr) -> bool:
    if cnvpr is None:
        return False
    for node in cnvpr.iter():
        if _local(node.tag) == "decorative":
            return str(node.get("val", "1")).lower() in ("1", "true")
    return False


# --- PowerPoint --------------------------------------------------------------

# Shape types that carry meaning visually and therefore need alt text.
_PPTX_NEEDS_ALT = {"PICTURE", "CHART", "DIAGRAM", "MEDIA", "GROUP", "LINKED_PICTURE"}


def check_pptx(path: str) -> list[dict]:
    from pptx import Presentation  # imported lazily so HTML/MD need no deps
    from pptx.util import Emu

    prs = Presentation(path)
    findings: list[dict] = []
    titles: dict[str, list[int]] = {}

    for index, slide in enumerate(prs.slides, start=1):
        where = f"Slide {index}"

        title_shape = slide.shapes.title
        title = title_shape.text.strip() if title_shape is not None else ""
        if not title:
            findings.append(finding(
                "Error", "All slides have titles", where,
                "The slide has no title placeholder, or its title is empty.",
                "Add a title in the slide's title placeholder. To keep it off the "
                "canvas, move it off-slide or set it via Accessibility > Slide Title.",
            ))
        else:
            titles.setdefault(title.lower(), []).append(index)

        positioned = []
        for shape in slide.shapes:
            shape_name = getattr(shape, "name", "shape")
            shape_where = f"{where} / {shape_name}"
            cnvpr = _find_local(shape._element, "cNvPr")
            alt = (cnvpr.get("descr") or "").strip() if cnvpr is not None else ""
            decorative = _is_decorative(cnvpr)
            kind = shape.shape_type.name if shape.shape_type is not None else ""

            has_text = bool(getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip())
            needs_alt = kind in _PPTX_NEEDS_ALT or (not has_text and kind not in ("TABLE",))

            if needs_alt and not alt and not decorative:
                findings.append(finding(
                    "Error", "All non-text content has alternative text (alt text)",
                    shape_where,
                    f"{kind or 'Shape'} has no alt text and is not marked decorative.",
                    "Describe what the object conveys in one sentence, or mark it "
                    "decorative if it carries no information.",
                ))

            if kind == "MEDIA":
                findings.append(finding(
                    "Warning", "Closed captions are included for inserted audio and video",
                    shape_where,
                    "Embedded media found. Captions cannot be verified from the file.",
                    "Confirm the media has closed captions or an accompanying transcript.",
                ))

            if getattr(shape, "has_table", False):
                table = shape.table
                if not table.first_row:
                    findings.append(finding(
                        "Error", "Tables specify column header information", shape_where,
                        "The table has no header row set.",
                        "Tick Header Row under Table Design so the first row is announced "
                        "as column headers.",
                    ))
                merged = any(
                    cell.is_merge_origin
                    for row in table.rows for cell in row.cells
                )
                if merged:
                    findings.append(finding(
                        "Warning", "Table has a simple structure", shape_where,
                        "The table contains merged cells.",
                        "Split merged cells, or restructure into separate simple tables.",
                    ))

            findings.extend(_pptx_contrast(shape, shape_where))

            top = shape.top if shape.top is not None else Emu(0)
            left = shape.left if shape.left is not None else Emu(0)
            positioned.append((int(top), int(left), shape_name))

        findings.extend(_reading_order(positioned, where))

    for title, slides in titles.items():
        if len(slides) > 1:
            findings.append(finding(
                "Tip", "Slide titles in a deck are unique",
                ", ".join(f"Slide {n}" for n in slides),
                f"{len(slides)} slides share the title {title!r}.",
                "Make each title distinct, e.g. by adding '(1 of 3)'.",
            ))

    return findings


def _pptx_contrast(shape, where: str) -> list[dict]:
    """Contrast for text drawn on an explicitly-filled shape.

    Ceiling: only explicit RGB values are compared. Theme colours, gradients,
    picture fills and inherited slide backgrounds are skipped rather than
    guessed at, the same blind spot Office's own checker has.
    """
    from pptx.dml.color import MSO_THEME_COLOR  # noqa: F401  (kept for clarity)
    from pptx.enum.dml import MSO_FILL

    if not getattr(shape, "has_text_frame", False):
        return []
    try:
        if shape.fill.type != MSO_FILL.SOLID:
            return []
        bg = shape.fill.fore_color
        if bg.type is None or not hasattr(bg, "rgb"):
            return []
        background = tuple(bg.rgb)  # type: ignore[arg-type]
    except (AttributeError, ValueError, TypeError):
        return []

    findings: list[dict] = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text.strip():
                continue
            try:
                colour = run.font.color
                if colour.type is None or not hasattr(colour, "rgb"):
                    continue
                foreground = tuple(colour.rgb)
            except (AttributeError, ValueError, TypeError):
                continue

            size = run.font.size.pt if run.font.size is not None else None
            large = size is not None and (size >= 18 or (size >= 14 and run.font.bold))
            threshold = AA_LARGE if large else AA_NORMAL
            ratio = contrast_ratio(foreground, background)  # type: ignore[arg-type]
            if ratio < threshold:
                findings.append(finding(
                    "Warning", "Sufficient contrast between text and background", where,
                    f"{run.text.strip()[:40]!r} is {round(ratio, 2)}:1 against its fill "
                    f"(WCAG AA needs {threshold}:1).",
                    "Darken the text or lighten the fill until the ratio is met.",
                ))
    return findings


def _reading_order(positioned: list[tuple[int, int, str]], where: str) -> list[dict]:
    """Compare z-order (what a screen reader follows) with visual top-left order.

    Ceiling: a geometry heuristic. Deliberate multi-column layouts can differ
    legitimately, so this is a Warning asking for a look, never an Error.
    """
    if len(positioned) < 3:
        return []
    row_tolerance = 250_000  # EMU, roughly a quarter inch
    expected = sorted(positioned, key=lambda s: (s[0] // row_tolerance, s[1]))
    if [s[2] for s in expected] == [s[2] for s in positioned]:
        return []
    return [finding(
        "Warning", "The reading order of the objects on a slide presentation is logical",
        where,
        "Screen-reader order is "
        + " → ".join(s[2] for s in positioned)
        + "; visual order looks like "
        + " → ".join(s[2] for s in expected)
        + ".",
        "Check Home > Arrange > Selection Pane and reorder, or confirm the current "
        "order is intentional.",
    )]


# --- Word --------------------------------------------------------------------


def check_docx(path: str) -> list[dict]:
    import docx

    document = docx.Document(path)
    findings: list[dict] = []

    for index, shape in enumerate(document.inline_shapes, start=1):
        doc_pr = _find_local(shape._inline, "docPr")
        alt = (doc_pr.get("descr") or "").strip() if doc_pr is not None else ""
        if not alt and not _is_decorative(doc_pr):
            name = (doc_pr.get("name") if doc_pr is not None else None) or f"image {index}"
            findings.append(finding(
                "Error", "All non-text content has alternative text (alt text)",
                f"Inline image: {name}",
                "The image has no alt text and is not marked decorative.",
                "Right-click > View Alt Text and describe what the image conveys.",
            ))

    if not any(p.style is not None and p.style.name.startswith("Heading") for p in document.paragraphs):
        findings.append(finding(
            "Tip", "Documents use heading styles", "Document",
            "No built-in Heading styles are used anywhere in the document.",
            "Apply Heading 1/2/3 styles instead of manually bolding and resizing text.",
        ))

    for index, table in enumerate(document.tables, start=1):
        where = f"Table {index}"
        # tblHeader must be set specifically on the first row to mean
        # anything; a mark on a later row doesn't repeat as a header row and
        # shouldn't count as satisfying this rule.
        first_row_header = table.rows and _find_local(table.rows[0]._tr, "tblHeader") is not None
        if not first_row_header:
            findings.append(finding(
                "Error", "Tables specify column header information", where,
                "The first row is not marked as a repeating header row.",
                "Select the first row, then Table Layout > Repeat Header Rows.",
            ))
        if _find_local(table._tbl, "gridSpan") is not None or _find_local(table._tbl, "vMerge") is not None:
            findings.append(finding(
                "Warning", "Table has a simple structure", where,
                "The table contains merged or split cells.",
                "Flatten the table so every row has the same number of cells.",
            ))

    return findings


# --- HTML --------------------------------------------------------------------


class _HTMLScan(HTMLParser):
    def __init__(self) -> None:
        super().__init__()
        self.findings: list[dict] = []
        self.lang: str | None = None
        self.headings: list[int] = []
        self._link_text: list[str] | None = None
        self._link_target = ""
        self._table_depth = 0
        self._table_has_th: list[bool] = []
        self._heading_level: int | None = None

    def handle_starttag(self, tag: str, attrs: list[tuple[str, str | None]]) -> None:
        attributes = {k.lower(): (v or "") for k, v in attrs}
        line = f"line {self.getpos()[0]}"

        if tag == "html":
            self.lang = attributes.get("lang") or None
        elif tag == "img":
            if "alt" not in attributes and attributes.get("role") != "presentation":
                self.findings.append(finding(
                    "Error", "All non-text content has alternative text (alt text)", line,
                    f"<img src={attributes.get('src', '')[:60]!r}> has no alt attribute.",
                    'Add alt="…", or alt="" if the image is purely decorative.',
                ))
        elif tag == "a":
            self._link_text = []
            self._link_target = attributes.get("href", "")
        elif tag == "table":
            self._table_depth += 1
            self._table_has_th.append(False)
        elif tag == "th" and self._table_has_th:
            self._table_has_th[-1] = True
        elif tag in ("h1", "h2", "h3", "h4", "h5", "h6"):
            self._heading_level = int(tag[1])
            self.headings.append(self._heading_level)

        style = attributes.get("style", "")
        if style:
            self.findings.extend(_style_contrast(style, line))

    def handle_endtag(self, tag: str) -> None:
        line = f"line {self.getpos()[0]}"
        if tag == "a" and self._link_text is not None:
            text = " ".join("".join(self._link_text).split()).strip().lower()
            if text in VAGUE_LINK_TEXT or (text.startswith("http") and len(text) > 30):
                self.findings.append(finding(
                    "Warning", "Hyperlink text is meaningful", line,
                    f"Link text {text[:40]!r} does not describe its destination.",
                    "Rewrite the link text to name the destination, e.g. "
                    "'Accessibility Checker rules'.",
                ))
            self._link_text = None
        elif tag == "table" and self._table_depth:
            self._table_depth -= 1
            if self._table_has_th and not self._table_has_th.pop():
                self.findings.append(finding(
                    "Error", "Tables specify column header information", line,
                    "The table has no <th> cells.",
                    "Mark the header row cells as <th scope=\"col\">.",
                ))
        elif tag.startswith("h") and self._heading_level is not None:
            self._heading_level = None

    def handle_data(self, data: str) -> None:
        if self._link_text is not None:
            self._link_text.append(data)


def _style_contrast(style: str, where: str) -> list[dict]:
    declarations = dict(
        (part.split(":", 1)[0].strip().lower(), part.split(":", 1)[1].strip())
        for part in style.split(";")
        if ":" in part
    )
    foreground = parse_color(declarations.get("color", ""))
    background = parse_color(declarations.get("background-color", "") or declarations.get("background", ""))
    if not foreground or not background:
        return []
    ratio = contrast_ratio(foreground, background)
    if ratio >= AA_NORMAL:
        return []
    return [finding(
        "Warning", "Sufficient contrast between text and background", where,
        f"Inline style gives {round(ratio, 2)}:1 (WCAG AA needs {AA_NORMAL}:1 for normal text).",
        "Adjust the colour pair until it reaches 4.5:1.",
    )]


def check_html(path: str) -> list[dict]:
    with open(path, encoding="utf-8", errors="replace") as handle:
        scanner = _HTMLScan()
        scanner.feed(handle.read())

    findings = scanner.findings
    if not scanner.lang:
        findings.append(finding(
            "Tip", "Document language is set", "<html>",
            "The <html> element has no lang attribute.",
            'Add lang="en" (or the document\'s actual language) so screen readers '
            "use the right pronunciation.",
        ))
    findings.extend(_heading_jumps(scanner.headings))
    return findings


def _heading_jumps(levels: list[int]) -> list[dict]:
    findings = []
    previous = 0
    for level in levels:
        if previous and level > previous + 1:
            findings.append(finding(
                "Tip", "Documents use heading styles", f"Heading level {level}",
                f"Heading level jumps from h{previous} to h{level}.",
                f"Use h{previous + 1} instead so the outline has no gaps.",
            ))
        previous = level
    return findings


# --- Markdown ----------------------------------------------------------------

_MD_IMAGE = re.compile(r"!\[(?P<alt>[^\]]*)\]\((?P<src>[^)\s]+)")
_MD_LINK = re.compile(r"(?<!!)\[(?P<text>[^\]]+)\]\((?P<href>[^)\s]+)")
_MD_HEADING = re.compile(r"^(?P<hashes>#{1,6})\s+\S")


def check_markdown(path: str) -> list[dict]:
    findings: list[dict] = []
    levels: list[int] = []

    with open(path, encoding="utf-8", errors="replace") as handle:
        for number, line in enumerate(handle, start=1):
            where = f"line {number}"
            for match in _MD_IMAGE.finditer(line):
                if not match.group("alt").strip():
                    # Empty alt (![](path)) is the standard Markdown/HTML
                    # convention for "this image is decorative", not
                    # automatically a failure the way a missing attribute
                    # would be in pptx/docx, so ask rather than demand text.
                    findings.append(finding(
                        "Tip", "All non-text content has alternative text (alt text)",
                        where,
                        f"Image {match.group('src')[:60]!r} has empty alt text.",
                        "Confirm this image is purely decorative. If it conveys "
                        "information, add a description: ![what it shows](path).",
                    ))
            for match in _MD_LINK.finditer(line):
                text = match.group("text").strip().lower()
                if text in VAGUE_LINK_TEXT:
                    findings.append(finding(
                        "Warning", "Hyperlink text is meaningful", where,
                        f"Link text {text!r} does not describe its destination.",
                        "Name the destination in the link text instead.",
                    ))
            heading = _MD_HEADING.match(line)
            if heading:
                levels.append(len(heading.group("hashes")))

    findings.extend(_heading_jumps(levels))
    return findings


# --- entry point -------------------------------------------------------------

CHECKERS = {
    ".pptx": check_pptx,
    ".docx": check_docx,
    ".html": check_html,
    ".htm": check_html,
    ".md": check_markdown,
    ".markdown": check_markdown,
}

SEVERITY_ORDER = {"Error": 0, "Warning": 1, "Tip": 2}


def run(path: str) -> list[dict]:
    extension = os.path.splitext(path)[1].lower()
    checker = CHECKERS.get(extension)
    if checker is None:
        raise SystemExit(
            f"Unsupported file type {extension!r}. Supported: {', '.join(sorted(CHECKERS))}"
        )
    try:
        findings = checker(path)
    except ImportError as exc:
        if extension == ".pptx":
            raise SystemExit(
                "python-pptx is required for .pptx files. Install it with: python -m pip install python-pptx"
            ) from exc
        if extension == ".docx":
            raise SystemExit(
                "python-docx is required for .docx files. Install it with: python -m pip install python-docx"
            ) from exc
        raise SystemExit(f"Missing Python dependency for {extension} files: {exc}") from exc

    findings.sort(key=lambda f: (SEVERITY_ORDER.get(f["severity"], 9), f["location"]))
    return findings


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("file", help="the .pptx, .docx, .html or .md file to check")
    parser.add_argument("--json", action="store_true", help="emit JSON instead of text")
    args = parser.parse_args()

    if not os.path.isfile(args.file):
        raise SystemExit(f"No such file: {args.file}")

    findings = run(args.file)

    if args.json:
        print(json.dumps({"file": args.file, "findings": findings}, indent=2))
        return 0

    if not findings:
        print(f"No mechanical accessibility issues found in {args.file}.")
        print("Alt-text accuracy, caption correctness and reading-order sense still need review.")
        return 0

    counts = {level: sum(1 for f in findings if f["severity"] == level) for level in SEVERITY_ORDER}
    print(f"{args.file}: {counts['Error']} errors, {counts['Warning']} warnings, {counts['Tip']} tips\n")
    for item in findings:
        print(f"[{item['severity']}] {item['rule']}")
        print(f"  Where: {item['location']}")
        print(f"  What:  {item['detail']}")
        print(f"  Fix:   {item['fix']}\n")
    return 0


if __name__ == "__main__":
    sys.exit(main())
