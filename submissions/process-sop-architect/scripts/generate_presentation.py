from __future__ import annotations

import argparse
import json
import math
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

try:
    from validate_process_spec import validate
except Exception:
    validate = None


def rgb(hex_value: str) -> RGBColor:
    hex_value = hex_value.replace('#', '')
    return RGBColor(int(hex_value[0:2], 16), int(hex_value[2:4], 16), int(hex_value[4:6], 16))


DEFAULT_THEME = {
    'colors': {
        'background': 'F7F9FC',
        'ink': '172033',
        'muted': '5F6B7A',
        'primary': '155EEF',
        'primary_dark': '0B3B8C',
        'secondary': '14B8A6',
        'accent': 'F97316',
        'success': '16A34A',
        'warning': 'F59E0B',
        'danger': 'DC2626',
        'line': 'D7DEE8',
        'white': 'FFFFFF',
    },
    'fonts': {'title': 'Aptos Display', 'body': 'Aptos'},
}


def load_theme(theme_path: Path | None) -> dict:
    if theme_path and theme_path.exists():
        data = json.loads(theme_path.read_text(encoding='utf-8'))
        merged = DEFAULT_THEME.copy()
        merged['colors'] = {**DEFAULT_THEME['colors'], **data.get('colors', {})}
        merged['fonts'] = {**DEFAULT_THEME['fonts'], **data.get('fonts', {})}
        return merged
    return DEFAULT_THEME


def add_bg(slide, colors):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = rgb(colors['background'])


def add_footer(slide, spec, slide_no, colors, fonts):
    tx = slide.shapes.add_textbox(Inches(0.45), Inches(7.05), Inches(8.0), Inches(0.25))
    p = tx.text_frame.paragraphs[0]
    p.text = f"{spec.get('process_name', 'Process')} | {spec.get('state', 'current').title()} State"
    p.font.size = Pt(8)
    p.font.name = fonts['body']
    p.font.color.rgb = rgb(colors['muted'])
    tx2 = slide.shapes.add_textbox(Inches(12.25), Inches(7.05), Inches(0.7), Inches(0.25))
    p2 = tx2.text_frame.paragraphs[0]
    p2.text = str(slide_no)
    p2.alignment = PP_ALIGN.RIGHT
    p2.font.size = Pt(8)
    p2.font.name = fonts['body']
    p2.font.color.rgb = rgb(colors['muted'])


def add_title(slide, title, subtitle, colors, fonts):
    box = slide.shapes.add_textbox(Inches(0.55), Inches(0.35), Inches(11.8), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = fonts['title']
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = rgb(colors['ink'])
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.58), Inches(0.92), Inches(11.8), Inches(0.34))
        sp = sub.text_frame.paragraphs[0]
        sp.text = subtitle
        sp.font.name = fonts['body']
        sp.font.size = Pt(12)
        sp.font.color.rgb = rgb(colors['muted'])


def add_round_rect(slide, x, y, w, h, fill, line=None, radius=True):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    if line:
        shp.line.color.rgb = rgb(line)
        shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    return shp


def add_text(slide, text, x, y, w, h, colors, fonts, size=14, bold=False, color='ink', align=None):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.text = str(text) if text is not None else ''
    p.font.name = fonts['body']
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = rgb(colors[color])
    if align:
        p.alignment = align
    return tb


def bullet_list(slide, items, x, y, w, h, colors, fonts, size=14, color='ink', max_items=6):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    for i, item in enumerate((items or [])[:max_items]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = str(item)
        p.level = 0
        p.font.name = fonts['body']
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(colors[color])
    if not items:
        p = tf.paragraphs[0]
        p.text = 'TBD'
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(colors['muted'])
    return tb


def kpi_card(slide, label, value, x, y, w, colors, fonts, accent='primary'):
    add_round_rect(slide, x, y, w, 1.05, colors['white'], colors['line'])
    add_text(slide, str(value), x + 0.18, y + 0.12, w - 0.36, 0.42, colors, fonts, size=24, bold=True, color=accent)
    add_text(slide, label, x + 0.18, y + 0.62, w - 0.36, 0.24, colors, fonts, size=10, color='muted')


def add_section_label(slide, label, x, y, colors, fonts, color='primary'):
    add_round_rect(slide, x, y, 1.55, 0.28, colors[color], None)
    add_text(slide, label.upper(), x + 0.1, y + 0.06, 1.35, 0.14, colors, fonts, size=7, bold=True, color='white', align=PP_ALIGN.CENTER)


def slide_cover(prs, spec, colors, fonts):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, colors)
    add_round_rect(slide, 0, 0, 13.333, 7.5, colors['primary_dark'], None, radius=False)
    add_round_rect(slide, 8.8, -0.4, 5.2, 8.4, colors['primary'], None, radius=False)
    add_text(slide, 'PROCESS & SOP ARCHITECT', 0.75, 0.7, 4.2, 0.3, colors, fonts, size=10, bold=True, color='white')
    title = spec.get('process_name', 'Process')
    add_text(slide, f'{title}\nProcess Briefing', 0.72, 2.05, 7.7, 1.45, colors, fonts, size=34, bold=True, color='white')
    subtitle = f"{spec.get('state', 'Current').title()} State | Version {spec.get('version','1.0')}"
    add_text(slide, subtitle, 0.78, 3.72, 6.8, 0.35, colors, fonts, size=17, color='white')
    purpose = spec.get('purpose', 'Purpose not stated')
    add_round_rect(slide, 0.78, 4.65, 7.35, 1.15, colors['white'], None)
    add_text(slide, purpose, 1.02, 4.89, 6.8, 0.66, colors, fonts, size=17, bold=False, color='ink')
    add_text(slide, 'Generated from validated process specification', 0.82, 6.55, 7.8, 0.25, colors, fonts, size=9, color='white')


def slide_snapshot(prs, spec, colors, fonts, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, colors); add_title(slide, 'The process is now packaged for operating alignment', 'Snapshot of scope, controls, risks, metrics, and unresolved questions.', colors, fonts); add_footer(slide, spec, num, colors, fonts)
    kpi_card(slide, 'Process steps', len(spec.get('steps', [])), 0.7, 1.55, 2.35, colors, fonts, 'primary')
    kpi_card(slide, 'Roles involved', len(spec.get('roles', [])), 3.25, 1.55, 2.35, colors, fonts, 'secondary')
    kpi_card(slide, 'Controls identified', len(spec.get('controls', [])), 5.8, 1.55, 2.35, colors, fonts, 'success')
    kpi_card(slide, 'Open questions', len(spec.get('open_questions', [])), 8.35, 1.55, 2.35, colors, fonts, 'warning')
    kpi_card(slide, 'Improvements', len(spec.get('improvements', [])), 10.9, 1.55, 1.75, colors, fonts, 'accent')
    scope = spec.get('scope', {})
    add_section_label(slide, 'Boundary', 0.7, 3.05, colors, fonts)
    add_round_rect(slide, 0.7, 3.45, 5.9, 1.7, colors['white'], colors['line'])
    add_text(slide, 'Starts', 0.95, 3.7, 1.05, 0.23, colors, fonts, size=10, bold=True, color='muted')
    add_text(slide, scope.get('starts', 'TBD'), 2.0, 3.65, 4.15, 0.35, colors, fonts, size=14)
    add_text(slide, 'Ends', 0.95, 4.35, 1.05, 0.23, colors, fonts, size=10, bold=True, color='muted')
    add_text(slide, scope.get('ends', 'TBD'), 2.0, 4.3, 4.15, 0.35, colors, fonts, size=14)
    add_section_label(slide, 'Inputs / Outputs', 7.0, 3.05, colors, fonts, 'secondary')
    add_round_rect(slide, 7.0, 3.45, 5.55, 1.7, colors['white'], colors['line'])
    bullet_list(slide, [f"Input: {x}" for x in spec.get('inputs', [])] + [f"Output: {x}" for x in spec.get('outputs', [])], 7.25, 3.68, 5.0, 1.05, colors, fonts, size=11, max_items=5)
    add_section_label(slide, 'Metrics', 0.7, 5.55, colors, fonts, 'accent')
    metrics = [f"{m.get('name')}: {m.get('target','TBD')}" for m in spec.get('metrics', [])]
    bullet_list(slide, metrics, 0.9, 5.95, 11.7, 0.75, colors, fonts, size=14, max_items=3)


def slide_process_flow(prs, spec, colors, fonts, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, colors); add_title(slide, 'The workflow clarifies handoffs, decisions, and rework loops', 'Process steps are grouped by owner and connected through documented paths.', colors, fonts); add_footer(slide, spec, num, colors, fonts)
    steps = spec.get('steps', [])[:10]
    owners = []
    for s in steps:
        if s.get('owner') not in owners:
            owners.append(s.get('owner'))
    left, top, width, height = 0.5, 1.45, 12.25, 5.25
    lane_h = height / max(len(owners), 1)
    positions = {}
    for i, owner in enumerate(owners):
        y = top + i * lane_h
        add_round_rect(slide, left, y, width, lane_h - 0.05, colors['white'], colors['line'], radius=False)
        add_round_rect(slide, left, y, 1.45, lane_h - 0.05, colors['primary_dark'] if i % 2 == 0 else colors['primary'], None, radius=False)
        add_text(slide, owner or 'TBD', left + 0.1, y + 0.18, 1.25, lane_h - 0.2, colors, fonts, size=10, bold=True, color='white')
    owner_counts = {o: 0 for o in owners}
    owner_steps = {o: [s for s in steps if s.get('owner') == o] for o in owners}
    max_count = max([len(v) for v in owner_steps.values()] or [1])
    box_w = min(1.65, (width - 1.8) / max(max_count, 1) - 0.15)
    for s in steps:
        owner = s.get('owner')
        lane_index = owners.index(owner) if owner in owners else 0
        order = owner_counts[owner]
        owner_counts[owner] += 1
        x = left + 1.65 + order * ((width - 1.9) / max(max_count, 1))
        y = top + lane_index * lane_h + lane_h * 0.22
        h = lane_h * 0.55
        positions[s.get('id')] = (x, y, box_w, h)
    # Connectors first so they sit behind nodes.
    for s in steps:
        sx = positions.get(s.get('id'))
        if not sx:
            continue
        targets = []
        if s.get('type') == 'decision':
            targets = [s.get('yes_next'), s.get('no_next')]
        elif s.get('next'):
            targets = [s.get('next')]
        for tgt in targets:
            if tgt not in positions:
                continue
            x1, y1, w1, h1 = sx
            x2, y2, w2, h2 = positions[tgt]
            conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1 + w1), Inches(y1 + h1 / 2), Inches(x2), Inches(y2 + h2 / 2))
            conn.line.color.rgb = rgb(colors['muted'])
            conn.line.width = Pt(1.1)
    for s in steps:
        x, y, w, h = positions[s.get('id')]
        is_decision = s.get('type') == 'decision'
        shape_type = MSO_AUTO_SHAPE_TYPE.DIAMOND if is_decision else MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        shp = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid(); shp.fill.fore_color.rgb = rgb(colors['warning'] if is_decision else colors['white'])
        shp.line.color.rgb = rgb(colors['primary_dark'] if not is_decision else colors['accent'])
        shp.line.width = Pt(1.2)
        tf = shp.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]; p.text = f"{s.get('id')}\n{s.get('name')}"; p.alignment = PP_ALIGN.CENTER; p.font.size = Pt(8.5); p.font.bold = True; p.font.name = fonts['body']; p.font.color.rgb = rgb(colors['ink'])
    if len(spec.get('steps', [])) > 10:
        add_text(slide, f"Note: showing first 10 of {len(spec.get('steps', []))} process steps. Full procedure is in the SOP.", 0.7, 6.72, 10, 0.25, colors, fonts, size=9, color='muted')


def slide_roles(prs, spec, colors, fonts, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, colors); add_title(slide, 'Role clarity is the foundation for standard execution', 'Primary owners, systems, and handoffs are visible across the process.', colors, fonts); add_footer(slide, spec, num, colors, fonts)
    roles = spec.get('roles', [])[:5]
    x = 0.7
    y = 1.55
    for i, role in enumerate(roles):
        add_round_rect(slide, x, y + i*0.83, 5.25, 0.62, colors['white'], colors['line'])
        add_text(slide, role.get('name','TBD'), x+0.2, y+0.1+i*0.83, 1.35, 0.2, colors, fonts, size=12, bold=True, color='primary')
        add_text(slide, role.get('description',''), x+1.55, y+0.1+i*0.83, 3.75, 0.28, colors, fonts, size=10, color='ink')
    systems = spec.get('systems', []) or ['TBD']
    add_section_label(slide, 'Systems', 6.55, 1.55, colors, fonts, 'secondary')
    for i, sys in enumerate(systems[:6]):
        add_round_rect(slide, 6.55 + (i%2)*2.75, 1.98 + (i//2)*0.75, 2.35, 0.45, colors['secondary'] if i == 0 else colors['white'], colors['line'])
        add_text(slide, sys, 6.7 + (i%2)*2.75, 2.1 + (i//2)*0.75, 2.0, 0.16, colors, fonts, size=10, bold=(i==0), color='white' if i==0 else 'ink', align=PP_ALIGN.CENTER)
    add_section_label(slide, 'Handoff indicators', 6.55, 4.7, colors, fonts, 'accent')
    handoffs = []
    previous_owner = None
    for step in spec.get('steps', []):
        owner = step.get('owner')
        if previous_owner and owner and owner != previous_owner:
            handoffs.append(f"{previous_owner} → {owner}: {step.get('name')}")
        previous_owner = owner
    bullet_list(slide, handoffs[:5], 6.65, 5.1, 5.6, 1.3, colors, fonts, size=12, max_items=5)


def slide_controls_risks(prs, spec, colors, fonts, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, colors); add_title(slide, 'Controls and risks show where governance must be explicit', 'Controls are tied to evidence; unresolved risk ownership remains visible.', colors, fonts); add_footer(slide, spec, num, colors, fonts)
    controls = spec.get('controls', [])[:5]
    add_section_label(slide, 'Controls', 0.65, 1.45, colors, fonts, 'success')
    for i, c in enumerate(controls):
        y = 1.85 + i*0.72
        add_round_rect(slide, 0.65, y, 5.95, 0.52, colors['white'], colors['line'])
        add_text(slide, c.get('id',''), 0.82, y+0.12, 0.42, 0.16, colors, fonts, size=10, bold=True, color='success')
        add_text(slide, c.get('name',''), 1.22, y+0.08, 2.05, 0.18, colors, fonts, size=10, bold=True)
        add_text(slide, f"Owner: {c.get('owner','TBD')} | Evidence: {c.get('evidence','TBD')}", 3.28, y+0.08, 3.0, 0.18, colors, fonts, size=8.5, color='muted')
    risks = spec.get('risks', [])[:5]
    add_section_label(slide, 'Risks', 7.0, 1.45, colors, fonts, 'danger')
    for i, r in enumerate(risks):
        y = 1.85 + i*0.72
        add_round_rect(slide, 7.0, y, 5.5, 0.52, colors['white'], colors['line'])
        add_text(slide, r.get('id',''), 7.18, y+0.12, 0.42, 0.16, colors, fonts, size=10, bold=True, color='danger')
        add_text(slide, r.get('description',''), 7.55, y+0.08, 2.2, 0.18, colors, fonts, size=10, bold=True)
        add_text(slide, f"Impact: {r.get('impact','TBD')}", 9.75, y+0.08, 2.35, 0.18, colors, fonts, size=8.5, color='muted')
    if spec.get('open_questions'):
        add_round_rect(slide, 0.65, 6.0, 11.85, 0.58, colors['warning'], None)
        add_text(slide, 'Open governance question: ' + spec.get('open_questions', [''])[0], 0.92, 6.17, 11.2, 0.18, colors, fonts, size=12, bold=True, color='white')


def slide_improvements(prs, spec, colors, fonts, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, colors); add_title(slide, 'Priority improvements convert findings into an action path', 'Recommendations remain separate from the documented current-state process.', colors, fonts); add_footer(slide, spec, num, colors, fonts)
    improvements = spec.get('improvements', [])[:6]
    phases = ['Now', 'Next', 'Later']
    phase_colors = ['primary', 'secondary', 'accent']
    for i, phase in enumerate(phases):
        x = 0.7 + i*4.1
        add_round_rect(slide, x, 1.55, 3.65, 0.5, colors[phase_colors[i]], None)
        add_text(slide, phase, x+0.1, 1.72, 3.45, 0.16, colors, fonts, size=11, bold=True, color='white', align=PP_ALIGN.CENTER)
        add_round_rect(slide, x, 2.18, 3.65, 3.75, colors['white'], colors['line'])
    for i, item in enumerate(improvements):
        phase_idx = min(i // 2, 2)
        x = 0.9 + phase_idx*4.1
        y = 2.48 + (i % 2)*1.58
        add_text(slide, item.get('title','Untitled improvement'), x, y, 3.15, 0.26, colors, fonts, size=13, bold=True, color='ink')
        add_text(slide, f"Value: {item.get('value','TBD')} | Effort: {item.get('effort','TBD')} | Owner: {item.get('owner','TBD')}", x, y+0.35, 3.1, 0.28, colors, fonts, size=9.5, color='muted')
        add_text(slide, item.get('recommendation',''), x, y+0.72, 3.1, 0.42, colors, fonts, size=10.5, color='ink')
    if not improvements:
        add_text(slide, 'No improvement opportunities were documented in the supplied process specification.', 1.0, 2.7, 10.8, 0.5, colors, fonts, size=16, color='muted', align=PP_ALIGN.CENTER)


def slide_next_steps(prs, spec, colors, fonts, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, colors); add_title(slide, 'Next steps close gaps before rollout or audit use', 'Use this page to drive validation, ownership, and implementation planning.', colors, fonts); add_footer(slide, spec, num, colors, fonts)
    add_section_label(slide, 'Open questions', 0.75, 1.55, colors, fonts, 'warning')
    bullet_list(slide, spec.get('open_questions', []), 0.95, 1.95, 5.3, 2.05, colors, fonts, size=13, max_items=6)
    add_section_label(slide, 'Assumptions', 7.0, 1.55, colors, fonts, 'primary')
    bullet_list(slide, spec.get('assumptions', []), 7.15, 1.95, 5.1, 2.05, colors, fonts, size=13, max_items=6)
    next_steps = [
        'Validate process scope and boundaries with the accountable process owner.',
        'Confirm missing owners, controls, and evidence requirements.',
        'Review the SOP and process map with performers and approvers.',
        'Prioritize the improvement backlog and assign implementation owners.',
        'Refresh the process package after decisions are confirmed.'
    ]
    add_section_label(slide, 'Recommended actions', 0.75, 4.55, colors, fonts, 'secondary')
    bullet_list(slide, next_steps, 0.95, 4.95, 11.2, 1.35, colors, fonts, size=13, max_items=5)


def add_simple_chart_slide(prs, spec, colors, fonts, num):
    slide = prs.slides.add_slide(prs.slide_layouts[6]); add_bg(slide, colors); add_title(slide, 'Process complexity is concentrated in steps, controls, risks, and improvements', 'This editable chart provides a quick orientation for leaders.', colors, fonts); add_footer(slide, spec, num, colors, fonts)
    chart_data = CategoryChartData()
    chart_data.categories = ['Steps', 'Roles', 'Controls', 'Risks', 'Improvements', 'Questions']
    chart_data.add_series('Count', [len(spec.get('steps', [])), len(spec.get('roles', [])), len(spec.get('controls', [])), len(spec.get('risks', [])), len(spec.get('improvements', [])), len(spec.get('open_questions', []))])
    chart = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(1.05), Inches(1.75), Inches(11.2), Inches(4.55), chart_data).chart
    chart.has_legend = False
    chart.category_axis.tick_labels.font.size = Pt(10)
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.chart_title.has_text_frame = True
    chart.chart_title.text_frame.text = 'Process package inventory'
    chart.chart_title.text_frame.paragraphs[0].font.size = Pt(14)
    for series in chart.series:
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = rgb(colors['primary'])


def generate_presentation(spec: dict, output_path: Path, theme_path: Path | None = None) -> Path:
    if validate is not None:
        errors, _warnings = validate(spec)
        if errors:
            raise ValueError('Validation errors: ' + '; '.join(errors))
    theme = load_theme(theme_path)
    colors = theme['colors']; fonts = theme['fonts']
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide_cover(prs, spec, colors, fonts)
    slide_snapshot(prs, spec, colors, fonts, 2)
    slide_process_flow(prs, spec, colors, fonts, 3)
    slide_roles(prs, spec, colors, fonts, 4)
    slide_controls_risks(prs, spec, colors, fonts, 5)
    slide_improvements(prs, spec, colors, fonts, 6)
    add_simple_chart_slide(prs, spec, colors, fonts, 7)
    slide_next_steps(prs, spec, colors, fonts, 8)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument('--input', required=True, type=Path, help='Path to process_spec.json')
    parser.add_argument('--output-dir', required=True, type=Path, help='Directory for generated presentation')
    parser.add_argument('--theme', type=Path, default=None, help='Optional theme JSON')
    args = parser.parse_args()
    spec = json.loads(args.input.read_text(encoding='utf-8'))
    slug = ''.join(ch.lower() if ch.isalnum() else '_' for ch in spec['process_name']).strip('_')
    deck_path = args.output_dir / f'{slug}_executive_process_briefing.pptx'
    theme_path = args.theme
    if theme_path is None:
        candidate = Path(__file__).resolve().parents[1] / 'assets' / 'presentation_theme.json'
        theme_path = candidate if candidate.exists() else None
    generate_presentation(spec, deck_path, theme_path)
    print(f'Presentation written: {deck_path}')


if __name__ == '__main__':
    main()
